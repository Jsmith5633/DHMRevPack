"""
Revenue Strategy Packet Tool — Streamlit App
Single-file build: slides 1–14 + 365-day + ops forecast all in one file.
No external imports needed. Drop into repo root alongside driftwood_logo.png.

requirements.txt:
  streamlit>=1.28.0
  python-pptx>=1.0.0
  pandas>=1.5.0
  openpyxl>=3.0.0
  lxml>=4.9.0
"""

import streamlit as st
import pandas as pd
import io, re, copy
from datetime import date, timedelta
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData

# ─── Page config ─────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Revenue Strategy Packet Tool",
    page_icon="🏨",
    layout="centered",
)

st.markdown("""
<style>
  .main { background: #F4F6F8; }
  .stButton>button {
    background: #0D1B2A; color: #fff; border: none;
    padding: 0.6rem 2rem; border-radius: 6px; font-size: 1rem; width: 100%;
  }
  .stButton>button:hover { background: #0A7E8C; }
  .status-box {
    background: #fff; border-left: 4px solid #0A7E8C;
    padding: 1rem 1.2rem; border-radius: 4px; margin: 0.5rem 0;
    font-size: 0.9rem; color: #2C3E50;
  }
  .prop-card {
    background: #0D1B2A; color: #fff; border-radius: 8px;
    padding: 1rem 1.4rem; margin: 0.8rem 0;
  }
  .prop-card h4 { color: #12A8B8; margin:0 0 0.3rem 0; font-size:0.8rem;
                  letter-spacing:1px; text-transform:uppercase; }
  .prop-card p  { margin:0; font-size:1rem; }
</style>
""", unsafe_allow_html=True)

# ─── Color palette ────────────────────────────────────────────────────────────
def _rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

C = {
    "navy":      _rgb("0D1B2A"), "teal":      _rgb("0A7E8C"),
    "tealLt":    _rgb("12A8B8"), "gold":      _rgb("D4A843"),
    "white":     _rgb("FFFFFF"), "offWhite":  _rgb("F4F6F8"),
    "slate":     _rgb("4A6274"), "lightGray": _rgb("E8EDF0"),
    "midGray":   _rgb("8FA3B1"), "green":     _rgb("27AE60"),
    "red":       _rgb("E74C3C"), "orange":    _rgb("E67E22"),
    "darkGray":  _rgb("2C3E50"), "amber":     _rgb("BF360C"),
    "grn_dk":    _rgb("1B5E20"), "grn_lt":    _rgb("D0EDD4"),
    "amb_lt":    _rgb("FFE8B0"), "red_lt":    _rgb("FCCEC9"),
    "grn_hdr":   _rgb("1A5C3A"), "wknd_bg":   _rgb("E8EBF8"),
    "wknd_hdr":  _rgb("1A2E6E"),
}

MO      = ['','Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
MO_FULL = ['','January','February','March','April','May','June',
           'July','August','September','October','November','December']
DAYS_MO = [0,31,28,31,30,31,30,31,31,30,31,30,31]

# ─── Safe value helpers ───────────────────────────────────────────────────────
def _fv(v, default=0.0):
    try:
        s = str(v).strip()
        if s in ('nan','','None'): return default
        return float(s)
    except: return default

def _iv(v, default=0): return int(_fv(v, default))

def _sign(v):
    try: return "+" if float(v) > 0 else ""
    except: return ""

# ─── PPTX primitives ─────────────────────────────────────────────────────────
def _new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color or C['white']

def _rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def _txt(slide, text, x, y, w, h, size=10, bold=False, color=None,
         align='left', valign='middle', wrap=False, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tb.text_frame._txBody.bodyPr.set('anchor',
        {'top':'t','middle':'ctr','bottom':'b'}.get(valign,'ctr'))
    p = tf.paragraphs[0]
    p.alignment = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,
                   'right':PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.name = 'Calibri'
    if color: r.font.color.rgb = color
    return tb

def _hdr(slide, title, sub=None):
    _rect(slide, 0, 0, 13.33, 0.55, C['navy'])
    _txt(slide, title, 0.25, 0, 9, 0.55, size=14, bold=True,
         color=C['white'], valign='middle')
    if sub:
        _txt(slide, sub, 8.5, 0, 4.6, 0.55, size=9,
             color=C['tealLt'], align='right', valign='middle')

def _cell(cell, text, bg=None, fg=None, bold=False, size=9,
          align='center', wrap=False):
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    tf = cell.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = {'center':PP_ALIGN.CENTER,'left':PP_ALIGN.LEFT,
                   'right':PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
    for run in p.runs: run.text = ''
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = 'Calibri'
    if fg: r.font.color.rgb = fg

def _kpi_dark(slide, x, y, w, h, label, value, sub, accent):
    _rect(slide, x, y, w, h, C['navy'])
    _rect(slide, x, y, w, 0.06, accent)
    _txt(slide, label, x+0.15, y+0.12, w-0.25, 0.25, size=8, color=C['midGray'])
    _txt(slide, value, x+0.15, y+0.35, w-0.25, 0.55, size=22, bold=True, color=C['white'])
    _txt(slide, sub,   x+0.15, y+0.88, w-0.25, 0.25, size=8, color=accent)

# ─── STR YoY inline XML patch ────────────────────────────────────────────────
def _add_yoy_xml(cell, chg_str, is_occ=False, is_index=False, is_total=False):
    """Insert small colored YoY run before endParaRPr in cell XML."""
    try: v = float(chg_str)
    except: return
    if v == 0: return
    sign = "+" if v > 0 else ""
    yoy = f"({sign}{v:.1f})" if (is_occ or is_index) else f"({sign}{v:.1f}%)"
    col = ("7FFFC4" if v > 0 else "FFB3B3") if is_total else \
          ("27AE60" if v > 0 else "E74C3C")
    tc = cell._tc
    txb = tc.find(qn('a:txBody'))
    if txb is None: return
    paras = txb.findall(qn('a:p'))
    if not paras: return
    mp = paras[0]
    runs = mp.findall(qn('a:r'))
    if not runs: return
    orig = runs[0].find(qn('a:rPr'))
    r2 = etree.Element(qn('a:r'))
    rp = etree.SubElement(r2, qn('a:rPr'))
    rp.set('lang','en-US'); rp.set('sz','600'); rp.set('dirty','0')
    sf = etree.SubElement(rp, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr')); sc.set('val', col)
    if orig is not None:
        for tag in [qn('a:latin'),qn('a:ea'),qn('a:cs')]:
            ch = orig.find(tag)
            if ch is not None: rp.append(copy.deepcopy(ch))
    t2 = etree.SubElement(r2, qn('a:t')); t2.text = f" {yoy}"
    end = mp.find(qn('a:endParaRPr'))
    if end is not None: end.addprevious(r2)
    else: mp.append(r2)

# ─── Data extraction ──────────────────────────────────────────────────────────
def _load_xl(file_bytes):
    return pd.ExcelFile(io.BytesIO(file_bytes))

def _parse_total_rooms(xl):
    """Parse total rooms dynamically from '0 / 150' format in 90 Day Segments col2."""
    try:
        df = pd.read_excel(xl, sheet_name='90 Day Segments', header=None)
        for idx in range(6, min(10, df.shape[0])):
            v = str(df.iloc[idx][2])
            if '/' in v:
                return int(v.split('/')[1].strip())
    except: pass
    return 150  # fallback

def _parse_pickup_notes(xl):
    """
    Parse segment pickup data from the notes text in Pickup sheet row 45 (0-indexed).
    Returns list of dicts: {name, rms, rev, pct_trans, is_total, note}
    """
    try:
        df = pd.read_excel(xl, sheet_name='Pickup', header=None)
        notes_text = str(df.iloc[45][1])
        if notes_text in ('nan', '', 'None'):
            return None

        # Find the first month's 7day pickup block (first occurrence)
        blocks = re.split(r'\n\n+', notes_text.strip())
        if not blocks:
            return None

        first_block = blocks[0]
        lines = [l.strip() for l in first_block.split('\n') if l.strip().startswith('*')]

        segs = []
        total_trn_rms = 0
        total_trn_rev = 0

        seg_map = {
            'transient': 'Transient (Total)',
            'retail':    'Retail / Rack',
            'discount':  'Discount',
            'internet':  'Internet/OTA',
            'ota':       'Internet/OTA',
            'package':   'Packages',
            'group':     'Group',
        }

        parsed = {}
        for line in lines:
            line_low = line.lower()
            name = None
            for key, label in seg_map.items():
                if key in line_low:
                    name = label
                    break
            if name is None:
                continue

            rms_match = re.search(r'([+-]?\d+)\s*RNs?', line, re.IGNORECASE)
            rev_match = re.search(r'\$([0-9,]+(?:\.[0-9]+)?)[k]?\s*in\s*revenue', line, re.IGNORECASE)
            adr_match = re.search(r'([+-]\$[\d.]+)\s*in\s*ADR', line, re.IGNORECASE)

            rms = int(rms_match.group(1)) if rms_match else 0
            rev_str = rev_match.group(1).replace(',','') if rev_match else '0'
            # Handle "$71k" shorthand
            rev_k_match = re.search(r'\$(\d+)k\b', line, re.IGNORECASE)
            if rev_k_match:
                rev = int(rev_k_match.group(1)) * 1000
            else:
                rev = int(float(rev_str)) if rev_str != '0' else 0

            adr_note = adr_match.group(1) if adr_match else ''

            if name == 'Transient (Total)':
                total_trn_rms = rms
                total_trn_rev = rev
                parsed[name] = {'rms': rms, 'rev': rev, 'note': f"ADR {adr_note}" if adr_note else '', 'is_total': True}
            else:
                parsed[name] = {'rms': rms, 'rev': rev, 'note': '', 'is_total': False}

        if not parsed:
            return None

        # Build ordered list
        order = ['Transient (Total)', 'Retail / Rack', 'Discount', 'Internet/OTA', 'Packages', 'Group']
        result = []
        for seg_name in order:
            d = parsed.get(seg_name)
            if d:
                pct = f"{d['rms']/total_trn_rms*100:.1f}%" if (total_trn_rms and not d['is_total']) else \
                      ('100.0%' if d['is_total'] else '')
                # Determine note for sub-segments
                note = d['note']
                if not note and not d['is_total']:
                    if seg_name == 'Retail / Rack': note = 'Strong rate integrity'
                    elif seg_name == 'Discount':     note = 'Monitor rate dilution'
                    elif seg_name == 'Internet/OTA': note = 'Healthy contribution'
                    elif seg_name == 'Packages':     note = 'Continue to promote'
                    elif seg_name == 'Group':        note = 'Solid group contribution'
                result.append({
                    'name': seg_name, 'rms': d['rms'], 'rev': d['rev'],
                    'pct_trans': pct, 'is_total': d['is_total'], 'note': note,
                })
            else:
                result.append({
                    'name': seg_name, 'rms': 0, 'rev': 0,
                    'pct_trans': '', 'is_total': seg_name == 'Transient (Total)', 'note': '—',
                })
        return result
    except:
        return None

def extract_all_data(file_bytes):
    """
    Master data extraction. Returns single dict used by all slide builders.
    All row/col indices verified against DTMIADR_rev_pak_2026-04-16.xlsx.
    """
    xl = _load_xl(file_bytes)
    df_str = pd.read_excel(xl, sheet_name='STR Analysis',    header=None)
    df_as  = pd.read_excel(xl, sheet_name='Annual Summary',  header=None)
    df_90  = pd.read_excel(xl, sheet_name='90 Day Segments', header=None)
    df_pu  = pd.read_excel(xl, sheet_name='Pickup',          header=None)
    df_365 = pd.read_excel(xl, sheet_name='365 Day Outlook', header=None)

    # ── Property info ────────────────────────────────────────────────────────
    total_rooms = _parse_total_rooms(xl)
    prop_name   = str(df_str.iloc[1][1])
    date_range  = str(df_str.iloc[4][1])   # "4/5/2026 - 4/11/2026"
    start_str   = date_range.split(' - ')[0].strip()
    parts       = start_str.split('/')
    report_mo   = int(parts[0])
    report_yr   = int(parts[2])
    end_str     = date_range.split(' - ')[1].strip()
    ep          = end_str.split('/')
    report_date = date(int(ep[2]), int(ep[0]), int(ep[1]))

    info = {
        'name':        prop_name,
        'date_range':  date_range,
        'report_mo':   report_mo,
        'report_yr':   report_yr,
        'report_date': report_date,
        'total_rooms': total_rooms,
    }

    # ── STR weekly + 28-day ──────────────────────────────────────────────────
    # DOW col pairs (val, chg): Sun(3,4) Mon(6,7) Tue(9,10) Wed(12,13)
    #                            Thu(15,16) Fri(18,19) Sat(21,22) Total(24,25)
    DOW_PAIRS = [(3,4),(6,7),(9,10),(12,13),(15,16),(18,19),(21,22),(24,25)]

    def _exrow(row_idx):
        row = df_str.iloc[row_idx]
        return [{'val': str(row[vc]) if str(row[vc]) not in ('nan','') else '',
                 'chg': str(row[cc]) if str(row[cc]) not in ('nan','') else ''}
                for vc, cc in DOW_PAIRS]

    # Weekly rows 8-18 (0-indexed), 28-day rows 24-34
    str_data = {
        'weekly': {
            'my_occ':    _exrow(8),  'cs_occ':    _exrow(9),  'mpi':       _exrow(10),
            'my_adr':    _exrow(12), 'cs_adr':    _exrow(13), 'ari':       _exrow(14),
            'my_revpar': _exrow(16), 'cs_revpar': _exrow(17), 'rgi':       _exrow(18),
        },
        'd28': {
            'my_occ':    _exrow(24), 'cs_occ':    _exrow(25), 'mpi':       _exrow(26),
            'my_adr':    _exrow(28), 'cs_adr':    _exrow(29), 'ari':       _exrow(30),
            'my_revpar': _exrow(32), 'cs_revpar': _exrow(33), 'rgi':       _exrow(34),
        },
    }

    # ── Week segment mix from STR Analysis ──────────────────────────────────
    # Segment names at row 37, data at row 46; block cols start at 12, step 4
    week_segs = []
    try:
        name_row  = df_str.iloc[37]
        total_row = df_str.iloc[46]
        seg_cols  = [ci for ci in range(12, df_str.shape[1], 4)
                     if str(name_row.iloc[ci]) not in ('nan','')]
        for nc in seg_cols:
            seg_name = str(name_row.iloc[nc]).strip()
            cy_rms   = _iv(total_row.iloc[nc])
            stly_rms = _iv(total_row.iloc[nc+1])
            cy_adr   = _fv(total_row.iloc[nc+2])
            stly_adr = _fv(total_row.iloc[nc+3])
            if cy_rms == 0 and stly_rms == 0: continue
            week_segs.append({
                'name': seg_name,
                'cy_rms': cy_rms, 'stly_rms': stly_rms,
                'cy_adr': cy_adr, 'stly_adr': stly_adr,
                'rms_var': cy_rms - stly_rms,
                'adr_var': cy_adr - stly_adr,
                'rev_var': cy_rms*cy_adr - stly_rms*stly_adr,
            })
    except: pass

    # ── Annual pace — OTB section rows 6-17 (0-indexed) ─────────────────────
    pace_data = []
    for i in range(6, 18):
        row = df_as.iloc[i]
        mo_name = str(row[1])
        if mo_name not in MO: continue
        mo_num = MO.index(mo_name)
        pace_data.append({
            'month': mo_name, 'mo_num': mo_num,
            'otb_occ':  _fv(row[2])*100,
            'otb_rms':  _iv(row[3]),
            'otb_adr':  _fv(row[4]),
            'otb_rev':  _fv(row[5]),
            'stly_occ': _fv(row[6])*100,
            'stly_rms': _iv(row[7]),
            'stly_adr': _fv(row[8]),
            'stly_rev': _fv(row[9]),
        })

    # ── Transient pace — rows 22-33 (0-indexed) ─────────────────────────────
    trn_map = {}
    for i in range(22, 34):
        row = df_as.iloc[i]
        mo_name = str(row[1])
        if mo_name not in MO: continue
        mo_num = MO.index(mo_name)
        trn_map[mo_num] = {
            'trn_rms':      _iv(row[5]),
            'trn_adr':      _fv(row[6]),
            'trn_stly_rms': _iv(row[8]),
            'trn_stly_adr': _fv(row[9]),
        }
    for d in pace_data:
        t = trn_map.get(d['mo_num'], {})
        d.update({
            'trn_rms':      t.get('trn_rms', 0),
            'trn_adr':      t.get('trn_adr', 0.0),
            'trn_stly_rms': t.get('trn_stly_rms', 0),
            'trn_stly_adr': t.get('trn_stly_adr', 0.0),
        })

    # ── BI Forecast + Budget — rows 38-49 (0-indexed) ───────────────────────
    fcst_map = {}
    for i in range(38, 50):
        row = df_as.iloc[i]
        mo_name = str(row[1])
        if mo_name not in MO: continue
        mo_num = MO.index(mo_name)
        fcst_map[mo_num] = {
            'bi_occ':  _fv(row[6]),
            'bi_rms':  _iv(row[7]),
            'bi_adr':  _fv(row[8]),
            'bi_rev':  _fv(row[9]),
            'bud_rms': _iv(row[21]),
            'bud_adr': _fv(row[22]),
            'bud_rev': _fv(row[23]),
        }
    for d in pace_data:
        f = fcst_map.get(d['mo_num'], {})
        d.update({
            'bi_occ':  f.get('bi_occ',  d['otb_occ']/100),
            'bi_rms':  f.get('bi_rms',  d['otb_rms']),
            'bi_adr':  f.get('bi_adr',  d['otb_adr']),
            'bi_rev':  f.get('bi_rev',  d['otb_rev']),
            'bud_rms': f.get('bud_rms', 0),
            'bud_adr': f.get('bud_adr', 0.0),
            'bud_rev': f.get('bud_rev', 0.0),
        })

    info['total_bi_rms']  = sum(d['bi_rms']  for d in pace_data)
    info['total_bi_rev']  = sum(d['bi_rev']  for d in pace_data)
    info['total_bud_rev'] = sum(d['bud_rev'] for d in pace_data)
    info['total_bi_adr']  = (info['total_bi_rev'] / info['total_bi_rms']
                              if info['total_bi_rms'] > 0 else 0)

    # ── 90 Day Segments: monthly total rows ──────────────────────────────────
    # Confirmed offsets: month+0 at row 36 (0-idx 35), +1 at 71 (70), +2 at 106 (105)
    MO_TOTAL_IDX = {0: 35, 1: 70, 2: 105}
    mo_totals = {}
    for offset, idx in MO_TOTAL_IDX.items():
        mo_num = report_mo + offset
        if mo_num > 12: break
        if idx >= df_90.shape[0]: break
        r = df_90.iloc[idx]
        mo_totals[mo_num] = {
            'lts':      _iv(r[3]),
            'occ':      _fv(r[4]),
            'adr':      _fv(r[5]),
            'rev':      _fv(r[6]),
            'trn_rms':  _iv(r[12]),
            'trn_adr':  _fv(r[13]),
            'grp_block':_iv(r[18]),
            'grp_rms':  _iv(r[19]),
        }

    # ── 90 Day daily data (OTB OCC) for monthly chart ───────────────────────
    # Block 0: rows 6-35 (0-idx), Block 1: 40-70, Block 2: 75-105
    BLOCK_RANGES = [(6,35), (40,70), (75,105)]
    mo_daily = {}
    for offset, (start, end) in enumerate(BLOCK_RANGES):
        mo_num = report_mo + offset
        if mo_num > 12: break
        days = []
        for i in range(start, min(end+1, df_90.shape[0])):
            row = df_90.iloc[i]
            dv = str(row[1])
            if '2026' not in dv and '2027' not in dv: continue
            try:
                pp = dv.split(',')[0].split('/')
                m2, d2 = int(pp[0]), int(pp[1])
                occ = _fv(row[4]) * 100
                days.append({'label': f"{m2}/{d2}", 'occ_otb': round(occ,1), 'occ_stly': 0})
            except: continue
        mo_daily[mo_num] = days

    # Overlay STLY OCC from 365 Day Outlook col 31 (STLY RMS / total_rooms)
    stly_map = {}
    for i in range(6, df_365.shape[0]):
        row = df_365.iloc[i]
        dv = str(row[1])
        if '2026' not in dv and '2027' not in dv: continue
        try:
            pp = dv.split(',')[0].split('/')
            m2, d2 = int(pp[0]), int(pp[1])
            stly_rms = _fv(row[31])
            stly_map[f"{m2}/{d2}"] = round(stly_rms / total_rooms * 100, 1)
        except: continue
    for days in mo_daily.values():
        for d in days:
            d['occ_stly'] = stly_map.get(d['label'], 0)

    # ── Pickup totals ─────────────────────────────────────────────────────────
    # Row 39 (0-idx): month labels at cols 2,4,6...
    # Row 40: RMS, Row 41: ADR, Row 42: REV
    pu_months, pu_rms, pu_adr, pu_rev = [], [], [], []
    pickup_from = ''
    try:
        mo_row  = df_pu.iloc[39]
        rms_row = df_pu.iloc[40]
        adr_row = df_pu.iloc[41]
        rev_row = df_pu.iloc[42]
        for ci in range(2, df_pu.shape[1], 2):
            mo_s = str(mo_row.iloc[ci])
            if mo_s in ('nan','','Total'): continue
            rms_v = _iv(rms_row.iloc[ci])
            adr_v = _fv(adr_row.iloc[ci])
            rev_v = _fv(rev_row.iloc[ci])
            if rms_v == 0 and adr_v == 0: continue
            pu_months.append(mo_s.split('-')[0])
            pu_rms.append(rms_v); pu_adr.append(round(adr_v)); pu_rev.append(rev_v)
    except: pass

    # Pickup "from" date — look in End Date sheet col header
    try:
        df_end = pd.read_excel(xl, sheet_name='End Date', header=None)
        for ci in range(df_end.shape[1]):
            v = str(df_end.iloc[4][ci])
            if 'Pickup From' in v:
                # "Pickup From 4/9/2026" → "Apr 9, 2026"
                m = re.search(r'(\d+)/(\d+)/(\d+)', v)
                if m:
                    mo2, d2, yr2 = int(m.group(1)), int(m.group(2)), int(m.group(3))
                    pickup_from = f"{MO[mo2]} {d2}, {yr2}"
                else:
                    pickup_from = v.split('From')[-1].strip()
                break
    except: pass

    info['pickup_from'] = pickup_from

    # ── Pickup segment data — parse from notes text ───────────────────────────
    pu_segs = _parse_pickup_notes(xl)
    if pu_segs is None:
        # Fallback: use totals only
        pu_segs = [
            {'name':'Transient (Total)', 'rms': pu_rms[0] if pu_rms else 0,
             'rev': int(pu_rev[0]) if pu_rev else 0, 'pct_trans': '100.0%',
             'is_total': True, 'note': f"ADR ${pu_adr[0]}" if pu_adr else ''},
        ]
        for seg in ['Retail / Rack','Discount','Internet/OTA','Packages','Group']:
            pu_segs.append({'name':seg,'rms':0,'rev':0,'pct_trans':'','is_total':False,'note':'—'})

    # ── 14-day ops forecast ───────────────────────────────────────────────────
    WEIGHTS = {
        1:(0.86,0.14), 2:(0.72,0.28), 3:(0.58,0.42), 4:(0.44,0.56),
        5:(0.30,0.70), 6:(0.16,0.84), 7:(0.00,1.00), 8:(0.00,1.14),
        9:(0.00,1.28),10:(0.00,1.42),11:(0.00,1.56),12:(0.00,1.70),
       13:(0.00,1.64),14:(0.00,0.00),
    }
    ops_rows = []
    for idx in range(6, df_90.shape[0]):
        row = df_90.iloc[idx]
        dv = str(row[1])
        if '2026' not in dv and '2027' not in dv: continue
        try:
            pp = dv.split(',')[0].split('/')
            mo2, d2, yr2 = int(pp[0]), int(pp[1]), int(pp[2])
            day_date = date(yr2, mo2, d2)
            delta = (day_date - report_date).days
            if delta < 1 or delta > 14: continue
        except: continue

        lts  = _iv(row[3])
        otb  = total_rooms - lts
        occ  = _fv(row[4])
        adr  = _fv(row[5])
        pu1d = _iv(row[7]); pu7d = _iv(row[9])
        trn  = _iv(row[12]); grp = _iv(row[19]); ctr = _iv(row[25])
        w1, w7 = WEIGHTS.get(delta, (0,0))
        pu_fcst    = round(pu1d*w1 + pu7d*w7)
        fcst_total = min(otb + pu_fcst, total_rooms)
        fcst_trn   = max(min(trn + pu_fcst, total_rooms - grp - ctr), trn)

        ops_rows.append({
            'date': f"{mo2}/{d2}", 'dow': day_date.strftime('%a'), 'dta': delta,
            'otb': otb, 'occ_pct': f"{occ*100:.1f}%",
            'adr': f"${adr:.2f}", 'pu1d': pu1d, 'pu7d': pu7d,
            'transient_otb': trn, 'group_otb': grp, 'contract_otb': ctr,
            'pu_fcst': pu_fcst, 'fcst_trn': fcst_trn, 'fcst_total': fcst_total,
            'fcst_occ': f"{fcst_total/total_rooms*100:.1f}%",
        })

    # ── 365-day data ──────────────────────────────────────────────────────────
    rows_365, comp_labels = _extract_365(df_365, report_date)

    return {
        'info':        info,
        'total_rooms': total_rooms,
        'str_data':    str_data,
        'week_segs':   week_segs,
        'pace_data':   pace_data,
        'mo_totals':   mo_totals,
        'mo_daily':    mo_daily,
        'pu_months':   pu_months,
        'pu_rms':      pu_rms,
        'pu_adr':      pu_adr,
        'pu_rev':      pu_rev,
        'pu_segs':     pu_segs,
        'ops_rows':    ops_rows,
        'rows_365':    rows_365,
        'comp_labels': comp_labels,
    }

def _smart_comp_abbrev(name):
    name = str(name).split('..')[0].strip()
    mappings = {
        'intercontinental':'IC','courtyard':'CY','aloft':'Aloft',
        'hilton garden':'HGI','residence inn':'RI','marriott':'Marr',
        'hyatt':'Hyatt','sheraton':'Sher','westin':'West',
        'embassy':'ES','homewood':'HW','hampton':'Hamp',
        'doubletree':'DT','springhill':'SHS','fairfield':'FF',
        'ac hotel':'AC','renaissance':'Rens','le meridien':'LM',
        'wyndham':'Wynd','days inn':'DI','quality inn':'QI',
        'holiday inn':'HI','best western':'BW','radisson':'Rad',
    }
    lower = name.lower()
    for key, val in mappings.items():
        if key in lower: return val
    return name[:5]

def _extract_365(df_365, report_date, max_comps=4):
    # Detect competitor columns (cols 23+)
    comp_cols = []
    for i in range(23, 33):
        try:
            name = str(df_365.iloc[5][i])
            if name == 'nan': continue
            val = df_365.iloc[6][i]
            if str(val) not in ('nan','0','0.0',''):
                comp_cols.append({'col':i, 'abbrev': _smart_comp_abbrev(name)})
                if len(comp_cols) >= max_comps: break
        except: pass

    def fv(v, dollar=False):
        try:
            f = float(v)
            if f == 0: return ''
            return f"${f:.0f}" if dollar else f"{f:.0f}"
        except: return ''

    rows = []
    for idx in range(6, df_365.shape[0]):
        row = df_365.iloc[idx]
        dv = str(row[1])
        if '2026' not in dv and '2027' not in dv: continue
        try:
            pp = dv.split(',')[0].split('/')
            mo2,d2,yr2 = int(pp[0]),int(pp[1]),int(pp[2])
            day_date = date(yr2,mo2,d2)
            dta = (day_date - report_date).days
        except: continue

        event = str(row[3]).strip() if str(row[3]) not in ('nan','') else ''
        if len(event) > 18: event = event[:16]+'..'

        grp_block = _fv(row[40]); grp_rms = _fv(row[41])
        blk_rem   = grp_block - grp_rms

        comp_vals = [fv(row[cc['col']]) for cc in comp_cols]
        while len(comp_vals) < max_comps: comp_vals.append('')

        rows.append({
            'date':      dv.split(',')[0],
            'dow':       day_date.strftime('%a'),
            'dta':       str(dta),
            'event':     event,
            'otb':       fv(row[5]),
            'occ':       f"{_fv(row[6]):.4f}",
            'lts':       fv(row[4]),
            'adr':       fv(row[7]),
            'hurdle':    fv(row[9]),
            'redeem':    fv(row[10]),
            'pu1d':      fv(row[12]),
            'pu1d_adr':  fv(row[13]),
            'pu7d':      fv(row[15]),
            'pu7d_adr':  fv(row[16]),
            'hotel_rate':fv(row[19], dollar=True),
            'avg_cs':    fv(row[21]),
            'comps':     comp_vals,
            'stly_var':  fv(row[34]),
            'trn_rms':   fv(row[37]),
            'trn_adr':   fv(row[38]),
            'grp_rms':   fv(row[41]),
            'blk_rem':   '' if blk_rem==0 else fv(blk_rem),
            'grp_adr':   fv(row[42]),
        })

    return rows, [c['abbrev'] for c in comp_cols]

# ─── Slide 1: Title / KPI Dashboard ──────────────────────────────────────────
def _build_slide_title(prs, data):
    info = data['info']
    w    = data['str_data']['weekly']
    TR   = data['total_rooms']
    slide = _blank(prs)
    _bg(slide, C['navy'])

    _txt(slide, "REVENUE STRATEGY",   0.5, 0.55, 12, 0.85, size=44, bold=True, color=C['white'])
    _txt(slide, "REVIEW & OUTLOOK",   0.5, 1.38, 12, 0.85, size=44, bold=True, color=C['tealLt'])
    _rect(slide, 0.5, 2.28, 1.2, 0.06, C['gold'])
    _txt(slide, info['name'],          0.5, 2.50, 12, 0.45, size=16, color=C['tealLt'])
    dr = info['date_range']
    _txt(slide, f"Week of {dr}  |  Revenue Package",
         0.5, 2.95, 12, 0.35, size=11, color=C['midGray'])

    def fv(key): return _fv(w[key][7]['val'])
    def fc(key): return w[key][7]['chg']

    # 6 KPI tiles: OCC, ADR, MPI, My RevPAR, CS RevPAR, RGI
    kpis = [
        ("OCC",       f"{fv('my_occ'):.1f}%",     f"CS: {fv('cs_occ'):.1f}%",
         f"{_sign(fc('my_occ'))}{_fv(fc('my_occ')):.1f} pts YoY", C['teal']),
        ("ADR",       f"${fv('my_adr'):.2f}",      f"CS: ${fv('cs_adr'):.2f}",
         f"{_sign(fc('my_adr'))}{_fv(fc('my_adr')):.1f}% YoY",    C['tealLt']),
        ("MPI",       f"{fv('mpi'):.1f}",           '',
         f"{_sign(fc('mpi'))}{_fv(fc('mpi')):.1f} pts YoY",        C['teal']),
        ("My RevPAR", f"${fv('my_revpar'):.2f}",    '',
         f"{_sign(fc('my_revpar'))}{_fv(fc('my_revpar')):.1f}% YoY", C['tealLt']),
        ("CS RevPAR", f"${fv('cs_revpar'):.2f}",    '',
         f"{_sign(fc('cs_revpar'))}{_fv(fc('cs_revpar')):.1f}% YoY", C['gold']),
        ("RGI",       f"{fv('rgi'):.1f}",            '',
         f"{_sign(fc('rgi'))}{_fv(fc('rgi')):.1f} pts YoY",
         C['green'] if fv('rgi') >= 100 else C['red']),
    ]
    W = 1.90; GAP = 0.20; START = 0.47; TOP = 5.10; H = 1.85
    for i, (lbl, val, cs_line, yoy_line, acc) in enumerate(kpis):
        left = START + i*(W+GAP)
        _rect(slide, left, TOP, W, H, C['navy'])
        _rect(slide, left, TOP, W, 0.07, acc)
        _txt(slide, lbl,      left, TOP+0.12, W, 0.28, size=11, bold=True, color=C['tealLt'], align='center')
        _txt(slide, val,      left, TOP+0.48, W, 0.72, size=22, bold=True, color=C['white'],  align='center')
        sub = cs_line if cs_line else yoy_line
        _txt(slide, sub,      left, TOP+1.18, W, 0.30, size=9,  color=C['midGray'], align='center')
        if cs_line:
            yoy_color = C['green'] if yoy_line.startswith('+') else C['red']
            _txt(slide, yoy_line, left, TOP+1.50, W, 0.25, size=9, color=yoy_color, align='center')

# ─── Slides 2–3: STR Weekly / 28-Day ─────────────────────────────────────────
def _build_str_slide(prs, data, period):
    info = data['info']
    w    = data['str_data'][period]
    segs = data['week_segs'] if period == 'weekly' else []

    slide = _blank(prs); _bg(slide)
    if period == 'weekly':
        title = f"STR PERFORMANCE — WEEK OF {info['date_range'].upper()}"
        sub   = f"{info['name'].split()[0]} vs Comp Set"
    else:
        title = "STR PERFORMANCE — RUNNING 28 DAYS (BY DAY OF WEEK)"
        sub   = f"{info['name'].split()[0]} vs Comp Set"
        _txt(slide, "28-DAY RUNNING SUMMARY", 0.25, 0.58, 8, 0.28, size=9, bold=True, color=C['teal'])
    _hdr(slide, title, sub)

    rows_def = [
        ("My OCC",    'my_occ',    True,  False, C['teal']),
        ("CS OCC",    'cs_occ',    True,  False, C['slate']),
        ("MPI",       'mpi',       False, True,  C['darkGray']),
        ("My ADR",    'my_adr',    False, False, C['teal']),
        ("CS ADR",    'cs_adr',    False, False, C['slate']),
        ("ARI",       'ari',       False, True,  C['darkGray']),
        ("My RevPAR", 'my_revpar', False, False, C['teal']),
        ("CS RevPAR", 'cs_revpar', False, False, C['slate']),
        ("RGI",       'rgi',       False, True,  C['darkGray']),
    ]

    tbl_y = 0.58 if period == 'weekly' else 0.88
    cw    = [1.55,1.4,1.4,1.4,1.4,1.4,1.4,1.4,1.45]
    tf = slide.shapes.add_table(10, 9, Inches(0.25), Inches(tbl_y), Inches(12.83), Inches(3.3))
    t  = tf.table
    for ci, w2 in enumerate(cw): t.columns[ci].width = int(w2*914400)
    for ri in range(10): t.rows[ri].height = int(0.33*914400)

    for ci, h in enumerate(["Metric","Sun","Mon","Tue","Wed","Thu","Fri","Sat","TOTAL"]):
        bg = _rgb("0A2438") if h == "TOTAL" else C['navy']
        fg = C['gold']      if h == "TOTAL" else C['white']
        _cell(t.cell(0,ci), h, bg=bg, fg=fg, bold=True, size=9.5)

    for ri, (lbl, key, is_occ, is_idx, color) in enumerate(rows_def):
        is_idx_row = key in ('mpi','ari','rgi')
        row_bg = C['teal'] if is_idx_row else (_rgb("E8EDF0") if ri%2==0 else C['white'])
        _cell(t.cell(ri+1,0), lbl, bg=row_bg,
              fg=C['white'] if is_idx_row else C['darkGray'], bold=is_idx_row, size=9.5)
        for ci, d in enumerate(w[key]):
            is_total = (ci == 7)
            bg = _rgb("0A2438") if is_total else (C['white'] if ri%2==0 else _rgb("F4F6F8"))
            val = d['val']
            if key in ('my_adr','cs_adr','my_revpar','cs_revpar'):
                try: val = f"${float(val):.0f}"
                except: pass
            _cell(t.cell(ri+1, ci+1), val, bg=bg,
                  fg=C['gold'] if is_total else color, bold=is_total, size=9.5)
            _add_yoy_xml(t.cell(ri+1, ci+1), d['chg'],
                         is_occ=is_occ, is_index=is_idx, is_total=is_total)

    if period == 'weekly':
        SEC_Y = 4.05
        _rect(slide, 0.25, SEC_Y, 7.2, 0.33, C['navy'])
        _txt(slide, "MARKET SEGMENT MIX — WEEK  |  CY vs STLY",
             0.25, SEC_Y, 7.2, 0.33, size=9.5, bold=True, color=C['white'], align='center', valign='middle')
        _rect(slide, 7.6, SEC_Y, 5.5, 3.4, _rgb("F4F6F8"))
        _rect(slide, 7.6, SEC_Y, 5.5, 0.33, C['navy'])
        _txt(slide, "STR ANALYSIS", 7.6, SEC_Y, 5.5, 0.33, size=10, bold=True,
             color=C['white'], align='center', valign='middle')

        rgi_val = w['rgi'][7]['val']; rgi_chg = w['rgi'][7]['chg']
        mpi_val = w['mpi'][7]['val']; ari_val = w['ari'][7]['val']
        my_adr_chg = w['my_adr'][7]['chg']; cs_adr_chg = w['cs_adr'][7]['chg']

        if segs:
            top_g = sorted(segs, key=lambda s: s['rms_var'], reverse=True)[:2]
            top_l = sorted(segs, key=lambda s: s['rms_var'])[:2]
            gainer_txt = " and ".join(
                f"{s['name']} growing {s['rms_var']:+d} rooms at "
                f"{_sign(str(s['adr_var']))}${s['adr_var']:.0f} ADR"
                for s in top_g if s['rms_var'] > 0) or "select segments"
            loser_txt = ", ".join(
                f"{s['name']} fell {abs(s['rms_var'])} rooms"
                for s in top_l if s['rms_var'] < 0) or ""
            seg_note = (f"The ARI improvement was driven by {gainer_txt} versus STLY. "
                        f"Offsetting this, {loser_txt}." if loser_txt else
                        f"The ARI improvement was driven by {gainer_txt} versus STLY.")
        else:
            seg_note = "Monitor segment mix and group pace to protect midweek production."

        critique = (
            f"The week of {info['date_range']} delivered an RGI of {rgi_val} "
            f"({_sign(rgi_chg)}{rgi_chg} pts YoY), "
            f"supported by an MPI of {mpi_val} and ARI of {ari_val}. "
            "The hotel outperformed the comp set on both occupancy and rate "
            f"simultaneously for the week. {seg_note} "
            f"Rate growth was driven by ADR improving "
            f"{_sign(my_adr_chg)}{my_adr_chg}% vs STLY while the comp set "
            f"grew {_sign(cs_adr_chg)}{cs_adr_chg}%."
        )
        _txt(slide, critique, 7.68, SEC_Y+0.42, 5.32, 2.88, size=9.5, color=C['darkGray'], wrap=True)

        n_rows = max(len(segs), 1)
        tf2 = slide.shapes.add_table(n_rows+1, 4, Inches(0.25), Inches(SEC_Y+0.33), Inches(7.2), Inches(3.07))
        t2 = tf2.table
        for ci, cw2 in enumerate([2.5,1.5,1.6,1.6]): t2.columns[ci].width = int(cw2*914400)
        for ri in range(n_rows+1): t2.rows[ri].height = int(0.27*914400)
        for ci, h in enumerate(["Segment","RMS Var","ADR Var","Rev Var"]):
            _cell(t2.cell(0,ci), h, bg=C['teal'], fg=C['white'], bold=True, size=9.5)
        if not segs:
            _cell(t2.cell(1,0), "No segment data in this export", bg=_rgb("F4F6F8"), fg=C['midGray'], size=9, align='left')
            for ci in range(1,4): _cell(t2.cell(1,ci),'—', bg=_rgb("F4F6F8"), fg=C['midGray'], size=9)
        else:
            for ri, seg in enumerate(segs):
                bg = _rgb("F4F6F8") if ri%2==0 else C['white']
                _cell(t2.cell(ri+1,0), seg['name'], bg=bg, fg=C['darkGray'], bold=True, size=9.5)
                _cell(t2.cell(ri+1,1), f"{seg['rms_var']:+d}", bg=bg,
                      fg=C['green'] if seg['rms_var']>0 else C['red'], size=9.5)
                _cell(t2.cell(ri+1,2), f"{_sign(str(seg['adr_var']))}${abs(seg['adr_var']):.0f}", bg=bg,
                      fg=C['green'] if seg['adr_var']>0 else C['orange'], size=9.5)
                _cell(t2.cell(ri+1,3), f"{_sign(str(seg['rev_var']))}${abs(seg['rev_var'])/1000:.1f}K", bg=bg,
                      fg=C['green'] if seg['rev_var']>0 else C['red'], size=9.5)
    else:
        rgi_tot = _fv(w['rgi'][7]['val']); mpi_tot = _fv(w['mpi'][7]['val'])
        ari_tot = _fv(w['ari'][7]['val']); rv_tot  = _fv(w['my_revpar'][7]['val'])
        rv_chg  = w['my_revpar'][7]['chg']; mpi_chg = w['mpi'][7]['chg']
        rgi_chg = w['rgi'][7]['chg'];       ari_chg = w['ari'][7]['chg']
        kpis = [
            ("28-Day MPI", f"{mpi_tot:.1f}", f"{_sign(mpi_chg)}{_fv(mpi_chg):.1f} pts", C['teal']),
            ("28-Day ARI", f"{ari_tot:.1f}", f"{_sign(ari_chg)}{_fv(ari_chg):.1f} pts", C['gold']),
            ("28-Day RGI", f"{rgi_tot:.1f}", f"{_sign(rgi_chg)}{_fv(rgi_chg):.1f} pts",
             C['green'] if rgi_tot>=100 else C['red']),
            ("My RevPAR",  f"${rv_tot:.2f}", f"{_sign(rv_chg)}{_fv(rv_chg):.1f}% YoY",  C['tealLt']),
        ]
        for i, (lbl, val, sub, acc) in enumerate(kpis):
            _kpi_dark(slide, 0.25+i*3.28, 4.2, 3.2, 1.3, lbl, val, sub, acc)

# ─── Slide 4: Annual Pace ─────────────────────────────────────────────────────
def _build_slide_annual_pace(prs, data):
    info = data['info']; pace = data['pace_data']
    slide = _blank(prs); _bg(slide)
    _hdr(slide, "ANNUAL PACE — OTB vs STLY vs BUDGET",
         f"{info['report_yr']} Full Year | Transient Focus")

    months   = [d['month'] for d in pace]
    otb_rms  = [d['otb_rms']  for d in pace]
    stly_rms = [d['stly_rms'] for d in pace]
    bud_rms  = [d['bud_rms']  for d in pace]

    cd = ChartData(); cd.categories = months
    cd.add_series(f"OTB {info['report_yr']}",   tuple(otb_rms))
    cd.add_series(f"STLY {info['report_yr']-1}", tuple(stly_rms))
    cd.add_series("Budget",                      tuple(bud_rms))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.25), Inches(0.65), Inches(13.0), Inches(3.8), cd).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = C['teal']
    chart.series[1].format.fill.solid(); chart.series[1].format.fill.fore_color.rgb = C['lightGray']
    chart.series[2].format.fill.solid(); chart.series[2].format.fill.fore_color.rgb = C['gold']

    upcoming = [d for d in pace if d['mo_num'] >= info['report_mo']][:3]
    headers  = ["Month","OTB RMS","OTB ADR","OTB REV","STLY RMS","STLY ADR","STLY REV",
                "RMS Var","ADR Var","REV Var","Bud REV","vs Bud"]
    cws = [0.7,0.9,0.9,0.9,0.9,0.9,0.9,0.7,0.7,0.9,0.9,0.7]
    tf = slide.shapes.add_table(len(upcoming)+1, len(headers),
         Inches(0.25), Inches(4.6), Inches(12.83), Inches(2.8))
    t = tf.table
    for ci, cw in enumerate(cws): t.columns[ci].width = int(cw*914400)
    for ri in range(len(upcoming)+1): t.rows[ri].height = int(0.55*914400)

    gc = {1:C['teal'],2:C['teal'],3:C['teal'],
          4:C['slate'],5:C['slate'],6:C['slate'],
          7:_rgb("1A5C3A"),8:_rgb("1A5C3A"),9:_rgb("1A5C3A"),
          10:C['gold'],11:_rgb("5C3A1A")}
    _cell(t.cell(0,0), "Month", bg=C['navy'], fg=C['white'], bold=True, size=8)
    for ci, h in enumerate(headers[1:], 1):
        _cell(t.cell(0,ci), h, bg=gc.get(ci,C['navy']), fg=C['white'], bold=True, size=7.5)

    for ri, d in enumerate(upcoming):
        bg = _rgb("F4F6F8") if ri%2==0 else C['white']
        rms_var = d['otb_rms']-d['stly_rms']; adr_var = d['otb_adr']-d['stly_adr']
        stly_rev = d['stly_rms']*d['stly_adr']; rev_var = d['otb_rev']-stly_rev
        vs_bud = (d['otb_rev']/d['bud_rev']-1)*100 if d['bud_rev']>0 else 0
        g=C['green']; r=C['red']
        vals = [
            (d['month'],             C['navy'],  True),
            (f"{d['otb_rms']:,}",    C['teal'],  False),
            (f"${d['otb_adr']:.2f}", C['teal'],  False),
            (f"${d['otb_rev']/1000:.0f}K", C['teal'], True),
            (f"{d['stly_rms']:,}",   C['slate'], False),
            (f"${d['stly_adr']:.2f}",C['slate'], False),
            (f"${stly_rev/1000:.0f}K",C['slate'],False),
            (f"{rms_var:+,}",        g if rms_var>=0 else r, True),
            (f"{_sign(str(adr_var))}${adr_var:.2f}", g if adr_var>=0 else r, False),
            (f"{_sign(str(rev_var))}${abs(rev_var)/1000:.0f}K", g if rev_var>=0 else r, True),
            (f"${d['bud_rev']/1000:.0f}K", C['midGray'], False),
            (f"{vs_bud:+.1f}%",      g if vs_bud>=0 else r, True),
        ]
        for ci, (v, fg, bold) in enumerate(vals):
            _cell(t.cell(ri+1,ci), v, bg=bg, fg=fg, bold=bold, size=8.5)

# ─── Slide 5: Transient Pace ──────────────────────────────────────────────────
def _build_slide_transient_pace(prs, data):
    info = data['info']; pace = data['pace_data']
    slide = _blank(prs); _bg(slide)
    _hdr(slide, "TRANSIENT PACE — OTB vs. STLY", "Transient Room Nights & ADR by Month")

    months   = [d['month'] for d in pace]
    trn_otb  = [d['trn_rms']      for d in pace]
    trn_stly = [d['trn_stly_rms'] for d in pace]
    otb_adr  = [d['trn_adr']      for d in pace]
    stly_adr = [d['trn_stly_adr'] for d in pace]

    cd = ChartData(); cd.categories = months
    cd.add_series(f"Transient OTB {info['report_yr']}",   tuple(trn_otb))
    cd.add_series(f"Transient STLY {info['report_yr']-1}", tuple(trn_stly))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.25), Inches(0.65), Inches(6.5), Inches(3.7), cd).chart
    chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = C['teal']
    chart.series[1].format.fill.solid(); chart.series[1].format.fill.fore_color.rgb = C['lightGray']

    cd2 = ChartData(); cd2.categories = months
    cd2.add_series(f"OTB ADR {info['report_yr']}",   tuple(otb_adr))
    cd2.add_series(f"STLY ADR {info['report_yr']-1}", tuple(stly_adr))
    chart2 = slide.shapes.add_chart(XL_CHART_TYPE.LINE,
        Inches(6.9), Inches(0.65), Inches(6.2), Inches(3.7), cd2).chart
    chart2.has_title = True; chart2.chart_title.text_frame.text = "ADR: OTB vs STLY"
    chart2.has_legend = True; chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart2.series[0].format.line.color.rgb = C['teal']; chart2.series[0].format.line.width = 18000
    chart2.series[1].format.line.color.rgb = C['gold']; chart2.series[1].format.line.width = 18000

    report_mo = info['report_mo']
    passed  = [d for d in pace if d['mo_num'] < report_mo]
    current = next((d for d in pace if d['mo_num'] == report_mo), None)
    future  = [d for d in pace if d['mo_num'] > report_mo][:2]

    def _tile(x, title, body, accent):
        _rect(slide, x, 4.55, 4.22, 2.5, C['white'])
        _rect(slide, x, 4.55, 4.22, 0.05, accent)
        _txt(slide, title, x+0.15, 4.6,  3.9, 0.3, size=9, bold=True, color=C['darkGray'])
        _txt(slide, body,  x+0.15, 4.95, 3.9, 2.0, size=9, color=C['slate'], wrap=True)

    if passed:
        yr_adr  = sum(d['trn_adr']      for d in passed)/len(passed)
        yr_sadr = sum(d['trn_stly_adr'] for d in passed)/len(passed)
        body = (f"OTB: {', '.join(str(d['trn_rms']) for d in passed[:3])}\n"
                f"STLY: {', '.join(str(d['trn_stly_rms']) for d in passed[:3])}\n"
                f"ADR +${yr_adr-yr_sadr:.2f} vs STLY avg")
        _tile(0.25, f"Jan–{passed[-1]['month']}", body, C['teal'])
    if current:
        var = current['trn_rms']-current['trn_stly_rms']
        denom = current['trn_stly_rms'] if current['trn_stly_rms'] else 1
        body = (f"OTB {current['trn_rms']:,} vs STLY {current['trn_stly_rms']:,}\n"
                f"{var:+,} RNs ({var/denom*100:+.1f}%)\n"
                f"ADR ${current['trn_adr']:.2f} vs ${current['trn_stly_adr']:.2f} STLY")
        _tile(4.56, f"{current['month']} OTB vs STLY", body, C['green'] if var>=0 else C['orange'])
    if future:
        lines = [f"OTB {d['trn_rms']:,} vs STLY {d['trn_stly_rms']:,} ({MO[d['mo_num']]})" for d in future]
        body = "\n".join(lines) + "\nMonitor forward pace closely"
        deficit = any(d['trn_rms']<d['trn_stly_rms'] for d in future)
        _tile(8.87, "May+ Pace Deficit" if deficit else "Forward Pace", body, C['orange'])

# ─── Slide 6: 7-Day Pickup ────────────────────────────────────────────────────
def _build_slide_pickup(prs, data):
    info = data['info']
    slide = _blank(prs); _bg(slide)
    _hdr(slide, "7-DAY PICKUP REPORT", f"Pickup from {info.get('pickup_from','')}")

    pu_months = data['pu_months']; pu_rms = data['pu_rms']
    pu_adr    = data['pu_adr'];    pu_rev = data['pu_rev']

    if pu_months and pu_rms:
        cd = ChartData(); cd.categories = pu_months
        cd.add_series("Room Nights Picked Up", tuple(pu_rms))
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.25), Inches(0.65), Inches(6.5), Inches(3.8), cd).chart
        chart.has_title = True; chart.chart_title.text_frame.text = "Room Nights Picked Up by Month"
        chart.has_legend = False
        chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = C['teal']
        chart.plots[0].has_data_labels = True; chart.plots[0].data_labels.show_value = True

    if pu_months and pu_adr:
        cd2 = ChartData(); cd2.categories = pu_months
        cd2.add_series("Average ADR of Pickup", tuple(pu_adr))
        chart2 = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(6.9), Inches(0.65), Inches(6.2), Inches(3.8), cd2).chart
        chart2.has_title = True; chart2.chart_title.text_frame.text = "Average ADR of Pickup by Month"
        chart2.has_legend = False
        chart2.series[0].format.fill.solid(); chart2.series[0].format.fill.fore_color.rgb = C['gold']
        chart2.plots[0].has_data_labels = True; chart2.plots[0].data_labels.show_value = True

    _rect(slide, 0.25, 4.62, 13.0, 0.33, C['navy'])
    _txt(slide, "PICKUP DETAIL (7-Day) — Key Transient Segments",
         0.35, 4.62, 13.0, 0.33, size=9.5, bold=True, color=C['white'], valign='middle')

    segs = data['pu_segs']
    tf = slide.shapes.add_table(len(segs)+1, 5,
        Inches(0.25), Inches(5.0), Inches(13.0), Inches(2.4))
    t = tf.table
    for ci, cw in enumerate([2.8,1.7,1.5,2.0,5.0]): t.columns[ci].width = int(cw*914400)
    for ri in range(len(segs)+1): t.rows[ri].height = int(0.36*914400)
    for ci, h in enumerate(["Segment","PU RNs","% of Trans","PU Revenue","Note"]):
        _cell(t.cell(0,ci), h, bg=C['teal'], fg=C['white'], bold=True, size=9.5)

    for ri, seg in enumerate(segs):
        is_total = seg.get('is_total', False)
        bg = _rgb("0D1B2A") if is_total else (_rgb("F4F6F8") if ri%2==0 else C['white'])
        fg = C['gold'] if is_total else C['darkGray']
        rms_v = seg['rms']
        rms_s = f"+{rms_v}" if rms_v>0 else (str(rms_v) if rms_v<0 else "—")
        rms_c = C['green'] if rms_v>0 else C['red'] if rms_v<0 else C['midGray']
        _cell(t.cell(ri+1,0), seg['name'], bg=bg, fg=fg, bold=is_total, size=9.5)
        _cell(t.cell(ri+1,1), rms_s,       bg=bg, fg=rms_c if not is_total else C['gold'], bold=is_total, size=9.5)
        _cell(t.cell(ri+1,2), seg.get('pct_trans',''), bg=bg, fg=C['teal'] if not is_total else C['gold'], size=9.5)
        _cell(t.cell(ri+1,3), f"${seg['rev']:,}" if seg['rev'] else '—', bg=bg,
              fg=C['teal'] if not is_total else C['gold'], bold=is_total, size=9.5)
        _cell(t.cell(ri+1,4), seg.get('note',''), bg=bg, fg=C['slate'] if not is_total else C['midGray'], size=9, align='left')

# ─── Slides 7/9/11: Monthly Daily Outlook ────────────────────────────────────
def _build_slide_monthly_occ(prs, data, mo_num):
    info     = data['info']; TR = data['total_rooms']
    mo_tot   = data['mo_totals'].get(mo_num, {})
    daily    = data['mo_daily'].get(mo_num, [])
    pace_row = next((d for d in data['pace_data'] if d['mo_num'] == mo_num), {})
    slide = _blank(prs); _bg(slide)
    mo_label = MO_FULL[mo_num] if mo_num <= 12 else ''
    _hdr(slide, f"{mo_label.upper()} DAILY OUTLOOK — OTB vs STLY",
         f"Transient  |  {mo_label} {info['report_yr']}")

    labels    = [d['label']             for d in daily]
    otb_vals  = [d['occ_otb']           for d in daily]
    stly_vals = [d.get('occ_stly', 0)   for d in daily]

    if labels:
        cd = ChartData(); cd.categories = labels
        cd.add_series(f"OTB {info['report_yr']} OCC%",    tuple(otb_vals))
        cd.add_series(f"STLY {info['report_yr']-1} OCC%", tuple(stly_vals))
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE,
            Inches(0.25), Inches(0.65), Inches(13.0), Inches(3.8), cd).chart
        chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.series[0].format.line.color.rgb = C['teal'];  chart.series[0].format.line.width = 18000
        chart.series[1].format.line.color.rgb = C['gold'];  chart.series[1].format.line.width = 18000
    else:
        _txt(slide, "No daily data available.", 0.25, 2.5, 13.0, 0.5, size=12, color=C['midGray'], align='center')

    otb_rms  = mo_tot.get('trn_rms',  0) or pace_row.get('otb_rms', 0)
    occ_pct  = (mo_tot.get('occ', 0) * 100) if mo_tot else pace_row.get('otb_occ', 0)
    trn_rms  = mo_tot.get('trn_rms',  pace_row.get('trn_rms', 0))
    trn_adr  = mo_tot.get('trn_adr',  pace_row.get('trn_adr', 0.0))
    stly_rms = pace_row.get('stly_rms', 0)
    vs_rms   = pace_row.get('otb_rms', 0) - stly_rms
    otb_rev  = pace_row.get('otb_rev', 0); stly_rev = pace_row.get('stly_rev', 0)
    vs_rev   = (otb_rev/stly_rev-1)*100 if stly_rev > 0 else 0
    vs_col   = C['green'] if vs_rms>=0 else C['red']
    rev_col  = C['green'] if vs_rev>=0 else C['red']

    kpis = [
        (f"{MO[mo_num]} OTB RMS",      f"{pace_row.get('otb_rms',0):,}", f"{occ_pct:.1f}% OCC", C['teal']),
        (f"{MO[mo_num]} Transient RMS", f"{trn_rms:,}",                   f"${trn_adr:.2f} ADR",  C['tealLt']),
        ("vs STLY Rooms",               f"{vs_rms:+,}",                   f"vs STLY {stly_rms:,}", vs_col),
        (f"{MO[mo_num]} Revenue OTB",   f"${otb_rev/1000:.0f}K",          f"{vs_rev:+.1f}% vs STLY", rev_col),
    ]
    for i, (lbl, val, sub, acc) in enumerate(kpis):
        _kpi_dark(slide, 0.25+i*3.26, 4.65, 3.2, 1.3, lbl, val, sub, acc)

# ─── Slides 8/10/12: Market Segment Mix ──────────────────────────────────────
def _build_slide_segment_mix(prs, data, mo_num):
    info  = data['info']
    segs  = data['week_segs']
    slide = _blank(prs); _bg(slide)
    mo_label = MO_FULL[mo_num] if mo_num <= 12 else ''
    _hdr(slide, f"MARKET SEGMENT MIX — {mo_label.upper()} {info['report_yr']}")

    if segs:
        total_rev = sum(max(s['cy_rms']*s['cy_adr'],0) for s in segs)
        if total_rev > 0:
            cd = ChartData()
            cd.categories = [s['name'] for s in segs]
            cd.add_series("Revenue %", tuple(max(s['cy_rms']*s['cy_adr'],0)/total_rev for s in segs))
            chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE,
                Inches(0.25), Inches(0.65), Inches(5.5), Inches(5.5), cd).chart
            chart.has_title = True; chart.chart_title.text_frame.text = "Revenue Contribution %"
            chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        tf = slide.shapes.add_table(len(segs)+1, 7,
            Inches(5.9), Inches(0.65), Inches(7.2), Inches(4.3))
        t = tf.table
        for ci, cw in enumerate([1.6,0.8,0.9,0.85,0.9,0.8,0.8]): t.columns[ci].width = int(cw*914400)
        for ri in range(len(segs)+1): t.rows[ri].height = int(0.52*914400)
        for ci, h in enumerate(["Segment","CY RNs","CY ADR","STLY RNs","STLY ADR","RN Var","ADR Var"]):
            _cell(t.cell(0,ci), h, bg=C['navy'], fg=C['white'], bold=True, size=9)
        for ri, seg in enumerate(segs):
            bg   = _rgb("F4F6F8") if ri%2==0 else C['white']
            rvar = seg['rms_var']; avar = seg['adr_var']
            rc   = C['green'] if rvar>0 else C['red'] if rvar<0 else C['midGray']
            ac   = C['green'] if avar>0 else C['orange'] if avar<0 else C['midGray']
            vals = [
                (seg['name'],             C['darkGray'], True),
                (f"{seg['cy_rms']:,}",    C['teal'],     False),
                (f"${seg['cy_adr']:.2f}", C['teal'],     True),
                (f"{seg['stly_rms']:,}",   C['slate'],   False),
                (f"${seg['stly_adr']:.2f}",C['slate'],   False),
                (f"{rvar:+,}" if rvar else "○", rc,     True),
                (f"{_sign(str(avar))}${avar:.2f}", ac,  False),
            ]
            for ci, (v, fg, bold) in enumerate(vals):
                _cell(t.cell(ri+1,ci), v, bg=bg, fg=fg, bold=bold, size=9)

        _rect(slide, 5.9, 5.1, 7.2, 0.32, C['teal'])
        _txt(slide, "KEY TAKEAWAYS", 6.0, 5.1, 7.0, 0.32, size=9, bold=True, color=C['white'], valign='middle')
        _rect(slide, 5.9, 5.42, 7.2, 2.0, _rgb("F4F6F8"))
        top2 = sorted(segs, key=lambda s: s['rms_var'], reverse=True)[:2]
        bot2 = sorted(segs, key=lambda s: s['rms_var'])[:2]
        lines = []
        for s in top2:
            if s['rms_var']>0: lines.append(f"✓ {s['name']} up {s['rms_var']:+,} RNs vs STLY at ${s['cy_adr']:.2f} ADR")
        for s in bot2:
            if s['rms_var']<0: lines.append(f"▲ {s['name']} decreased {s['rms_var']:,} RNs vs STLY — monitor")
        body = "\n".join(lines[:4]) if lines else "Review segment performance vs STLY."
        _txt(slide, body, 6.0, 5.48, 6.9, 1.85, size=9, color=C['darkGray'], wrap=True)
    else:
        _txt(slide, "Segment data not available in this export.",
             0.25, 3.5, 13.0, 0.5, size=11, color=C['midGray'], align='center')

# ─── Slide 13: Full Year Forecast ─────────────────────────────────────────────
def _build_slide_full_year(prs, data):
    info = data['info']; pace = data['pace_data']; TR = data['total_rooms']
    slide = _blank(prs); _bg(slide)
    _hdr(slide, f"FULL YEAR FORECAST {info['report_yr']}", "Finance Forecast vs Budget")

    months   = [d['month'] for d in pace]
    fcst_rms = [d['bi_rms'] for d in pace]
    bud_rms  = [d['bud_rms'] for d in pace]

    cd = ChartData(); cd.categories = months
    cd.add_series("Forecast", tuple(fcst_rms)); cd.add_series("Budget", tuple(bud_rms))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE,
        Inches(0.25), Inches(0.65), Inches(7.4), Inches(4.8), cd).chart
    chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.series[0].format.line.color.rgb = C['teal'];  chart.series[0].format.line.width = 20000
    chart.series[1].format.line.color.rgb = C['gold'];  chart.series[1].format.line.width = 20000

    fwd = [d for d in pace if d['mo_num'] >= info['report_mo']]
    tf = slide.shapes.add_table(len(fwd)+1, 7, Inches(7.65), Inches(0.65), Inches(5.45), Inches(4.8))
    t = tf.table
    for ci, cw in enumerate([0.75,0.85,0.85,0.9,0.85,0.75,0.5]): t.columns[ci].width = int(cw*914400)
    for ri in range(len(fwd)+1): t.rows[ri].height = int(0.38*914400)
    for ci, h in enumerate(["Month","OTB OCC","Fcst OCC","Fcst ADR","Fcst REV","Budget REV","vs Bud"]):
        _cell(t.cell(0,ci), h, bg=C['navy'], fg=C['white'], bold=True, size=8)
    for ri, d in enumerate(fwd):
        bg = _rgb("F4F6F8") if ri%2==0 else C['white']
        fcst_occ_p = d['bi_occ']*100 if d['bi_occ']<=1 else d['bi_occ']
        vs_bud = (d['bi_rev']/d['bud_rev']-1)*100 if d['bud_rev']>0 else 0
        vs_col = C['green'] if vs_bud>=0 else C['red']
        fcst_col = C['grn_dk'] if fcst_occ_p>=95 else C['amber']
        _cell(t.cell(ri+1,0), d['month'],                  bg=bg, fg=C['navy'],   bold=True, size=8)
        _cell(t.cell(ri+1,1), f"{d['otb_occ']:.1f}%",     bg=bg, fg=C['teal'] if d['otb_occ']>80 else C['orange'], size=8)
        _cell(t.cell(ri+1,2), f"{fcst_occ_p:.1f}%",        bg=bg, fg=fcst_col,   bold=True, size=8)
        _cell(t.cell(ri+1,3), f"${d['bi_adr']:.2f}",       bg=bg, fg=C['teal'],   size=8)
        _cell(t.cell(ri+1,4), f"${d['bi_rev']/1000:.0f}K", bg=bg, fg=C['navy'],   bold=True, size=8)
        _cell(t.cell(ri+1,5), f"${d['bud_rev']/1000:.0f}K",bg=bg, fg=C['midGray'],size=8)
        _cell(t.cell(ri+1,6), f"{vs_bud:+.1f}%",           bg=bg, fg=vs_col,      bold=True, size=8)

    total_bi_rev = info['total_bi_rev']; total_bud_rev = info['total_bud_rev']
    total_bi_rms = info['total_bi_rms']; total_bi_adr  = info['total_bi_adr']
    yr_occ = total_bi_rms / (365 * TR) * 100
    vs_bud_pct = (total_bi_rev/total_bud_rev-1)*100 if total_bud_rev>0 else 0

    kpis = [
        ("Forecast Full Year OCC", f"{yr_occ:.1f}%",            "", C['teal']),
        ("Forecast Full Year ADR", f"${total_bi_adr:.2f}",      "", C['gold']),
        ("Forecast Full Year REV", f"${total_bi_rev/1e6:.2f}M", "", C['tealLt']),
        ("vs Budget REV",          f"{vs_bud_pct:+.1f}%",       "",
         C['green'] if vs_bud_pct>=0 else C['red']),
    ]
    for i, (lbl, val, sub, acc) in enumerate(kpis):
        _kpi_dark(slide, 0.25+i*3.28, 5.6, 3.2, 1.3, lbl, val, sub, acc)

    commentary = (
        f"Finance Forecast projects full year at ${total_bi_rev/1e6:.2f}M vs "
        f"${total_bud_rev/1e6:.2f}M budget ({vs_bud_pct:+.1f}%). "
        "Heavy lift required in back half as current OTB reflects early pace. "
        "Transient volume capture is the primary revenue lever."
    )
    _rect(slide, 0.25, 7.05, 13.0, 0.38, _rgb("F0F4F8"))
    _rect(slide, 0.25, 7.05, 0.06, 0.38, C['teal'])
    _txt(slide, commentary, 0.4, 7.07, 12.7, 0.34, size=8.5, color=C['slate'], wrap=True)

# ─── Slide 14: 14-Day Operational Forecast ───────────────────────────────────
def _build_slide_ops(prs, data):
    rows14 = data['ops_rows']; TR = data['total_rooms']; info = data['info']
    if not rows14: return

    slide = _blank(prs); _bg(slide)
    _rect(slide, 0, 0, 13.33, 0.55, C['navy'])
    _txt(slide, "14-DAY OPERATIONAL FORECAST", 0.25, 0, 9, 0.55, size=14, bold=True, color=C['white'], valign='middle')
    start_d = rows14[0]['date']; end_d = rows14[-1]['date']
    _txt(slide, f"{start_d} – {end_d}, {info['report_yr']}  |  {TR} Total Rooms",
         7.0, 0, 6.0, 0.55, size=9, color=C['tealLt'], align='right', valign='middle')

    _rect(slide, 0.25, 0.62, 13.0, 0.28, _rgb("FFF8E7"))
    _txt(slide,
         f"Forecast = OTB + (1-Day PU × 1-Day Weight) + (7-Day PU × 7-Day Weight), "
         f"capped at {TR} rooms  |  Group = OTB fixed  |  Transient = OTB + all pickup",
         0.32, 0.62, 12.8, 0.28, size=7.5, color=C['amber'], wrap=False)

    n_cols = len(rows14) + 1
    tf = slide.shapes.add_table(13, n_cols, Inches(0.25), Inches(0.95), Inches(13.0), Inches(6.22))
    tbl = tf.table

    col_w = 12.65 / (n_cols - 1)
    tbl.columns[0].width = int(0.95 * 914400)
    for ci in range(1, n_cols):
        tbl.columns[ci].width = int(col_w * 914400)
    for ri in range(13):
        tbl.rows[ri].height = int(0.478 * 914400)

    def sc(cell, text, bg=None, fg=None, bold=False, fs=8.5):
        if bg: cell.fill.solid(); cell.fill.fore_color.rgb = bg
        tf2 = cell.text_frame; tf2.word_wrap = False
        p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        for r in p.runs: r.text = ''
        run = p.add_run(); run.text = str(text)
        run.font.size = Pt(fs); run.font.bold = bold; run.font.name = 'Calibri'
        if fg: run.font.color.rgb = fg

    def is_wknd(dow): return dow in ['Fri','Sat']

    def occ_bg(pct_str):
        try:
            v = float(pct_str.rstrip('%'))
            if v >= 95: return C['grn_lt'], C['grn_dk']
            if v >= 80: return C['amb_lt'], C['amber']
            return C['red_lt'], C['red']
        except: return C['white'], C['darkGray']

    ROW_LABELS = ['OTB OCC%','Total OTB Rms','Transient OTB','Contract OTB','Group OTB',
                  'Fcst PU Add','FCST TOTAL RMS','FCST OCC%','Fcst Transient','Fcst Contract',
                  'Fcst Group','OTB ADR']
    ROW_KEYS   = ['occ_pct','otb','transient_otb','contract_otb','group_otb',
                  'pu_fcst','fcst_total','fcst_occ','fcst_trn','contract_otb',
                  'group_otb','adr']

    # Header row
    sc(tbl.cell(0,0), '', bg=C['navy'])
    for ci, r in enumerate(rows14, 1):
        wknd = is_wknd(r['dow'])
        bg = C['wknd_hdr'] if wknd else C['navy']
        sc(tbl.cell(0,ci), f"{r['date']}\n{r['dow']}", bg=bg, fg=C['white'], bold=True, fs=8)

    # Data rows
    for ri, (label, key) in enumerate(zip(ROW_LABELS, ROW_KEYS)):
        is_fcst_total = (label == 'FCST TOTAL RMS')
        is_fcst_occ   = (label == 'FCST OCC%')
        is_otb_occ    = (label == 'OTB OCC%')
        lbl_bg = _rgb("1A5C3A") if 'FCST' in label else C['navy']
        sc(tbl.cell(ri+1, 0), label, bg=lbl_bg, fg=C['white'], bold=True, fs=8)

        for ci, r in enumerate(rows14, 1):
            wknd = is_wknd(r['dow'])
            val  = r.get(key, '')
            if key == 'pu_fcst':
                val = f"+{r['pu_fcst']}" if r['pu_fcst'] else '—'
            if key == 'group_otb' and label == 'Fcst Group':
                val = r['group_otb']
            if key == 'contract_otb' and label == 'Fcst Contract':
                val = r['contract_otb']

            if is_fcst_occ:
                cbg, cfg = occ_bg(str(val))
                sc(tbl.cell(ri+1,ci), val, bg=cbg, fg=cfg, bold=True, fs=8.5)
            elif is_otb_occ:
                cbg, cfg = occ_bg(str(val))
                sc(tbl.cell(ri+1,ci), val, bg=cbg, fg=cfg, fs=8.5)
            elif is_fcst_total:
                try:
                    v = int(val)
                    cbg = C['grn_lt'] if v >= int(TR*0.95) else C['amb_lt'] if v >= int(TR*0.80) else C['red_lt']
                    cfg = C['grn_dk'] if v >= int(TR*0.95) else C['amber']  if v >= int(TR*0.80) else C['red']
                except: cbg = C['white']; cfg = C['darkGray']
                sc(tbl.cell(ri+1,ci), val, bg=cbg, fg=cfg, bold=True, fs=8.5)
            else:
                if wknd: bg = C['wknd_bg']
                elif (ri+1) % 2 == 0: bg = C['offWhite']
                else: bg = C['white']
                fg = C['gold'] if (key in ('group_otb',) and val and val != '—') else C['slate']
                disp = str(val) if val and val != 0 else '—'
                if label in ('Group OTB','Fcst Group') and (not val or val == 0): disp = '—'
                sc(tbl.cell(ri+1,ci), disp, bg=bg, fg=fg, fs=8.5)

    # Legend
    leg_y = 7.28
    legends = [(_rgb("D0EDD4"),'Fcst ≥95%'),(_rgb("FFE8B0"),'Fcst 80–94%'),
               (_rgb("FCCEC9"),'Fcst <80%'),(C['wknd_hdr'],'Weekend col')]
    for i, (col, lbl) in enumerate(legends):
        lx = 0.3 + i*2.6
        _rect(slide, lx, leg_y+0.02, 0.16, 0.13, col)
        _txt(slide, lbl, lx+0.2, leg_y, 2.3, 0.2, size=7.5, color=C['slate'])

# ─── 365-Day slides ───────────────────────────────────────────────────────────
def _build_365_slides(prs, data):
    rows_365   = data['rows_365']
    comp_labels= data['comp_labels']
    report_date= data['info']['report_date']

    MO_NAMES = ['','Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
    ROWS_PER = 16; ROW_H_IN = 0.340; HDR_H = 0.52; BAND_H = 0.20; TABLE_X = 0.28

    COLS = [
        {'key':'date_dow','label':'Date',     'w':1.20},
        {'key':'dta',     'label':'DTA',      'w':0.38},
        {'key':'event',   'label':'Event',    'w':0.82},
        {'key':'otb',     'label':'OTB Rms',  'w':0.50},
        {'key':'occ',     'label':'OCC',      'w':0.44},
        {'key':'lts',     'label':'LTS',      'w':0.38},
        {'key':'adr',     'label':'OTB ADR',  'w':0.55},
        {'key':'hurdle',  'label':'Hurdle',   'w':0.50},
        {'key':'redeem',  'label':'Redm',     'w':0.38},
        {'key':'pu1d',    'label':'1D PU',    'w':0.40},
        {'key':'pu1d_adr','label':'1D ADR',   'w':0.48},
        {'key':'pu7d',    'label':'7D PU',    'w':0.40},
        {'key':'pu7d_adr','label':'7D ADR',   'w':0.48},
        {'key':'stly_var','label':'vs STLY',  'w':0.50},
        {'key':'trn_rms', 'label':'Trn Rms',  'w':0.50},
        {'key':'trn_adr', 'label':'Trn ADR',  'w':0.50},
        {'key':'grp_rms', 'label':'Grp Rms',  'w':0.46},
        {'key':'blk_rem', 'label':'Blk Rem',  'w':0.46},
        {'key':'grp_adr', 'label':'Grp ADR',  'w':0.46},
        {'key':'hotel_rate','label':'Hotel $','w':0.52},
        {'key':'avg_cs',  'label':'Avg CS',   'w':0.46},
    ]
    for i, lbl in enumerate(comp_labels[:4]):
        COLS.append({'key':f'comp{i}','label':lbl,'w':0.46})

    TABLE_W = sum(c['w'] for c in COLS)

    GROUPS = [
        {'label':'',             'keys':['date_dow','dta','event'],              'color':'0D1B2A'},
        {'label':'ON THE BOOKS', 'keys':['otb','occ','lts','adr','hurdle','redeem'],'color':'0A7E8C'},
        {'label':'PICKUP',       'keys':['pu1d','pu1d_adr','pu7d','pu7d_adr'],  'color':'1A5C6E'},
        {'label':'vs STLY',      'keys':['stly_var'],                             'color':'1A5C3A'},
        {'label':'TRANSIENT OTB','keys':['trn_rms','trn_adr'],                   'color':'4A6274'},
        {'label':'GROUP OTB',    'keys':['grp_rms','blk_rem','grp_adr'],         'color':'5C4A1A'},
        {'label':'RATE SHOP',    'keys':['hotel_rate','avg_cs']+[f'comp{i}' for i in range(len(comp_labels[:4]))], 'color':'1A3A5C'},
    ]

    def occ_style(occ_str):
        try:
            v = float(occ_str)
            if v >= 0.95: return C['grn_lt'], C['grn_dk']
            if v >= 0.70: return C['amb_lt'], C['amber']
            return C['red_lt'], C['red']
        except: return C['white'], C['darkGray']

    def stly_style(v_str):
        try:
            v = int(float(v_str))
            if v > 0:  return C['grn_lt'], C['grn_dk']
            if v < 0:  return C['red_lt'], C['red']
        except: pass
        return C['white'], C['darkGray']

    # Group rows into chunks of ROWS_PER
    chunks = [rows_365[i:i+ROWS_PER] for i in range(0, len(rows_365), ROWS_PER)]

    for chunk_idx, chunk in enumerate(chunks):
        if not chunk: continue

        slide = _blank(prs); _bg(slide)
        first_date = chunk[0]['date']; last_date = chunk[-1]['date']
        try:
            fp = first_date.split('/'); lp = last_date.split('/')
            mo1 = MO_NAMES[int(fp[0])]; yr1 = fp[2]
            mo2 = MO_NAMES[int(lp[0])]; yr2 = lp[2]
            hdr_mo = f"{mo1} {yr1}" if mo1==mo2 else f"{mo1}–{mo2} {yr1}"
        except: hdr_mo = ""

        _rect(slide, 0, 0, 13.33, 0.52, C['navy'])
        _txt(slide, f"365-DAY OUTLOOK — {hdr_mo.upper()}", 0.25, 0, 8, 0.52,
             size=13, bold=True, color=C['white'], valign='middle')
        _txt(slide, f"{first_date} – {last_date}", 8.5, 0, 4.6, 0.52,
             size=9, color=C['tealLt'], align='right', valign='middle')

        # Group header band
        band_y = HDR_H
        key_to_col_idx = {c['key']:i for i,c in enumerate(COLS)}
        for grp in GROUPS:
            indices = [key_to_col_idx[k] for k in grp['keys'] if k in key_to_col_idx]
            if not indices: continue
            grp_x = TABLE_X + sum(COLS[j]['w'] for j in range(min(indices)))
            grp_w = sum(COLS[j]['w'] for j in indices)
            _rect(slide, grp_x, band_y, grp_w, BAND_H, _rgb(grp['color']))
            if grp['label']:
                _txt(slide, grp['label'], grp_x, band_y, grp_w, BAND_H,
                     size=6, bold=True, color=C['white'], align='center', valign='middle')

        # Column header row
        col_hdr_y = band_y + BAND_H
        for ci, col in enumerate(COLS):
            cx = TABLE_X + sum(c['w'] for c in COLS[:ci])
            grp_color = '0D1B2A'
            for grp in GROUPS:
                if col['key'] in grp['keys']: grp_color = grp['color']; break
            _rect(slide, cx, col_hdr_y, col['w'], 0.26, _rgb(grp_color))
            _txt(slide, col['label'], cx, col_hdr_y, col['w'], 0.26,
                 size=6.5, bold=True, color=C['white'], align='center', valign='middle')

        # Data rows
        data_y = col_hdr_y + 0.26
        for ri, row in enumerate(chunk):
            ry = data_y + ri * ROW_H_IN
            dow = row.get('dow','')
            is_wknd = dow in ['Fri','Sat']
            row_bg = C['wknd_bg'] if is_wknd else (C['offWhite'] if ri%2==0 else C['white'])

            # Full-width background
            _rect(slide, TABLE_X, ry, TABLE_W, ROW_H_IN, row_bg)

            for ci, col in enumerate(COLS):
                cx = TABLE_X + sum(c['w'] for c in COLS[:ci])
                key = col['key']

                if key == 'date_dow':
                    val = f"{row['date']} {dow}"
                    _txt(slide, val, cx, ry, col['w'], ROW_H_IN, size=8, bold=is_wknd,
                         color=C['wknd_hdr'] if is_wknd else C['navy'], valign='middle')
                    continue

                if key.startswith('comp'):
                    idx = int(key[4:])
                    val = row['comps'][idx] if idx < len(row['comps']) else ''
                elif key == 'date_dow': val = ''
                else: val = row.get(key,'')

                if key == 'occ':
                    cbg, cfg = occ_style(val)
                    _rect(slide, cx, ry, col['w'], ROW_H_IN, cbg)
                    try: disp = f"{float(val)*100:.1f}"
                    except: disp = val
                    _txt(slide, disp, cx, ry, col['w'], ROW_H_IN, size=8, bold=True,
                         color=cfg, align='center', valign='middle')
                elif key == 'stly_var':
                    cbg, cfg = stly_style(val)
                    _rect(slide, cx, ry, col['w'], ROW_H_IN, cbg)
                    _txt(slide, val, cx, ry, col['w'], ROW_H_IN, size=8, bold=True,
                         color=cfg, align='center', valign='middle')
                elif key == 'hotel_rate':
                    _txt(slide, val, cx, ry, col['w'], ROW_H_IN, size=8,
                         color=C['navy'], align='center', valign='middle', bold=True)
                elif key == 'event':
                    _txt(slide, val, cx, ry, col['w'], ROW_H_IN, size=7,
                         color=C['amber'], align='left', valign='middle', italic=True, wrap=False)
                else:
                    _txt(slide, val, cx, ry, col['w'], ROW_H_IN, size=8,
                         color=C['darkGray'], align='center', valign='middle')

        # Legend
        leg_items = [
            (C['grn_lt'], C['grn_dk'], 'OCC ≥95%'),
            (C['amb_lt'], C['amber'],  'OCC 70–94%'),
            (C['red_lt'], C['red'],    'OCC <70%'),
            (C['wknd_bg'],C['wknd_hdr'],'Weekend'),
        ]
        ley = data_y + len(chunk)*ROW_H_IN + 0.05
        if ley < 7.2:
            for i, (bg, fg, lbl) in enumerate(leg_items):
                lx = TABLE_X + i*2.5
                _rect(slide, lx, ley, 0.18, 0.16, bg)
                _txt(slide, lbl, lx+0.22, ley, 2.2, 0.18, size=7.5, color=C['slate'])

# ─── Master builder ───────────────────────────────────────────────────────────
def build_presentation(file_bytes, progress_cb=None):
    def prog(pct, msg):
        if progress_cb: progress_cb(pct, msg)

    prog(0.05, "Reading Excel file...")
    data = extract_all_data(file_bytes)
    info = data['info']

    prog(0.20, "Building presentation...")
    prs = _new_prs()

    prog(0.25, "Slide 1: Title & KPI Dashboard...")
    _build_slide_title(prs, data)

    prog(0.30, "Slide 2: STR Weekly...")
    _build_str_slide(prs, data, 'weekly')

    prog(0.33, "Slide 3: STR 28-Day...")
    _build_str_slide(prs, data, 'd28')

    prog(0.36, "Slide 4: Annual Pace...")
    _build_slide_annual_pace(prs, data)

    prog(0.40, "Slide 5: Transient Pace...")
    _build_slide_transient_pace(prs, data)

    prog(0.44, "Slide 6: 7-Day Pickup...")
    _build_slide_pickup(prs, data)

    prog(0.48, "Slides 7–12: Monthly Outlook...")
    for offset in range(3):
        mo = info['report_mo'] + offset
        if mo > 12: break
        _build_slide_monthly_occ(prs, data, mo)
        _build_slide_segment_mix(prs, data, mo)

    prog(0.55, "Slide 13: Full Year Forecast...")
    _build_slide_full_year(prs, data)

    prog(0.60, "Slide 14: 14-Day Ops Forecast...")
    _build_slide_ops(prs, data)

    prog(0.65, "Slides 15–39: 365-Day Outlook...")
    _build_365_slides(prs, data)

    prog(0.95, "Saving...")
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    prog(1.00, "Done!")
    return buf.getvalue(), info

# ─── Streamlit UI ─────────────────────────────────────────────────────────────
def main():
    col_logo, col_title = st.columns([1,3])
    with col_logo:
        try:
            import base64
            with open("driftwood_logo.png","rb") as f:
                logo_b64 = base64.b64encode(f.read()).decode()
            st.markdown(
                f"<img src='data:image/png;base64,{logo_b64}' "
                f"style='width:100%;max-width:160px;padding-top:8px;'/>",
                unsafe_allow_html=True)
        except: pass
    with col_title:
        st.markdown("""
        <div style='padding-top:0.6rem'>
            <h1 style='margin:0;color:#0D1B2A;font-size:1.6rem;font-family:Calibri,sans-serif;'>
                Revenue Strategy Packet Tool
            </h1>
            <p style='margin:0;color:#0A7E8C;font-size:0.9rem;'>Driftwood Hospitality Management</p>
        </div>""", unsafe_allow_html=True)

    st.markdown("Upload your weekly rev pak Excel file to generate the revenue strategy presentation.")
    st.markdown("---")

    with st.expander("📋 How to prepare your rev pak file"):
        st.markdown("""
        1. Open **Lighthouse** and open the **RevPAK**
        2. In the RevPAK dropdown menu open **Annual Summary**
        3. Change all forecast types to **Finance** and all transient/group comparisons to **STLY**
        4. Export via **email**
        5. Once the email extract is received, upload here and generate file
        """)

    uploaded = st.file_uploader("Choose your rev pak file", type=['xlsx','xls'],
                                 help="The weekly revenue package Excel file from Lighthouse")

    if uploaded is None:
        st.info("👆 Upload your rev pak Excel file above to get started.")
        st.markdown("""
        **What this builds:**
        - Title slide with 6 KPI tiles (OCC, ADR, MPI, My RevPAR, CS RevPAR, RGI)
        - STR Performance (weekly + 28-day running)
        - 7-Day Pickup Report with segment detail
        - 14-Day Operational Forecast (Transient / Contract / Group split)
        - Monthly Outlook — current month + 2 forward
        - Full Year Forecast vs Budget
        - 365-Day Outlook (25 columns)
        """)
        return

    file_bytes = uploaded.read()
    st.success(f"✓ File loaded: **{uploaded.name}** ({len(file_bytes)/1024:.0f} KB)")

    # Preview
    try:
        xl   = _load_xl(file_bytes)
        df_s = pd.read_excel(xl, sheet_name='STR Analysis', header=None)
        prop_name  = str(df_s.iloc[1][1])
        date_range = str(df_s.iloc[4][1])
        tr = _parse_total_rooms(xl)
        parts = date_range.split(' - ')[0].split('/')
        report_mo = int(parts[0]); report_yr = int(parts[2])

        col1, col2 = st.columns(2)
        with col1:
            st.markdown(f"""<div class="prop-card"><h4>Property</h4><p>{prop_name}</p></div>""",
                        unsafe_allow_html=True)
        with col2:
            st.markdown(f"""<div class="prop-card"><h4>Report Period</h4><p>{date_range}</p></div>""",
                        unsafe_allow_html=True)
        col3, col4 = st.columns(2)
        with col3:
            st.markdown(f"""<div class="prop-card"><h4>Report Month</h4>
                <p>{MO[report_mo]} {report_yr}</p></div>""", unsafe_allow_html=True)
        with col4:
            st.markdown(f"""<div class="prop-card"><h4>Total Rooms</h4>
                <p>{tr} rooms</p></div>""", unsafe_allow_html=True)
    except Exception as e:
        st.warning(f"Could not preview file: {e}")

    st.markdown("---")

    if st.button("🚀 Generate Presentation", type="primary"):
        progress_bar = st.progress(0)
        status_text  = st.empty()

        def update_progress(pct, msg):
            progress_bar.progress(pct)
            status_text.markdown(f'<div class="status-box">⚙️ {msg}</div>', unsafe_allow_html=True)

        try:
            pptx_bytes, build_info = build_presentation(file_bytes, update_progress)
            progress_bar.progress(1.0)
            status_text.markdown(
                '<div class="status-box" style="border-color:#27AE60">✅ Presentation ready!</div>',
                unsafe_allow_html=True)

            prop_code = build_info['name'].split()[0].upper()
            mo_str    = MO[build_info['report_mo']]
            yr_str    = str(build_info['report_yr'])
            filename  = f"{prop_code}_Revenue_Strategy_{mo_str}{yr_str}.pptx"

            st.download_button(
                label="📥 Download PPTX",
                data=pptx_bytes,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except Exception as e:
            progress_bar.empty()
            st.error(f"❌ Build failed: {str(e)}")
            with st.expander("Error details"):
                import traceback
                st.code(traceback.format_exc())

    st.markdown("---")
    st.markdown(
        "<small style='color:#8FA3B1'>Revenue Strategy Packet Tool • Driftwood Hospitality Management • "
        "Upload a new file at any time to regenerate.</small>",
        unsafe_allow_html=True)

if __name__ == "__main__":
    main()
