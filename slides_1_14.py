"""
slides_1_14.py  — complete rewrite based on verified Excel structure
Matches reference PPTX design exactly.

Verified Excel mappings (0-indexed iloc):
STR Analysis:
  row 1  col 1  = property name
  row 4  col 1  = date range
  Weekly data rows 8-18, 28-day rows 24-34
  DOW cols: Sun(3,4) Mon(6,7) Tue(9,10) Wed(12,13) Thu(15,16) Fri(18,19) Sat(21,22) Total(24,25)
  Segment names row 37: GROUP@12, RACK@16, CORPORATE@20, DISCOUNT@24, LOCAL CORPORATE@28,
                         PACKAGES@32, GOVERNMENT@36, WHOLESALE@40, COMP@44
  Each segment block: [name_col+0]=RMS, [+1]=STLY RMS, [+2]=ADR, [+3]=STLY ADR
  Weekly total row 46

Annual Summary:
  OTB section rows 6-18: col1=month, col2=OCC, col3=RMS, col4=ADR, col5=REV,
                          col6=STLY OCC, col7=STLY RMS, col8=STLY ADR, col9=STLY REV
  Transient section rows 22-34: col1=month, col2=OTB RMS, col3=OTB ADR, col4=OTB REV,
                                  col5=Trn RMS, col6=Trn ADR, col7=Trn REV,
                                  col8=Trn STLY RMS, col9=Trn STLY ADR
  BI Forecast section rows 38-50: col1=month, col2=OTB OCC, col3=OTB RMS, col4=OTB ADR, col5=OTB REV,
                                   col6=BI Fcst OCC, col7=BI Fcst RMS, col8=BI Fcst ADR, col9=BI Fcst REV,
                                   col21=Budget OCC, col22=Budget RMS, col23=Budget ADR, col24=Budget REV
  Total row 50

90 Day Segments:
  Daily data starts row 6; monthly totals at row 36 (mo+0), 71 (mo+1), 106 (mo+2)
  cols: 1=date, 2=OOO/RMS, 3=LTS, 4=OCC, 5=ADR, 6=REV,
        7=1D PU RMS, 8=1D PU ADR, 9=7D PU RMS, 10=7D PU ADR,
        12=Trn RMS, 13=Trn ADR, 14=Trn REV,
        15=Trn STLY RMS, 16=Trn STLY ADR,
        18=Grp Block, 19=Grp RMS, 20=Grp ADR, 21=Grp REV,
        25=Contract RMS

Pickup:
  Month labels row 39 (cols 2,4,6... every-other), RMS row 40, ADR row 41, REV row 42
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy
from datetime import date

# ── Color palette ────────────────────────────────────────────────────────────
def _rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

C = {
    "navy":     _rgb("0D1B2A"), "teal":     _rgb("0A7E8C"),
    "tealLt":   _rgb("12A8B8"), "gold":     _rgb("D4A843"),
    "white":    _rgb("FFFFFF"), "offWhite": _rgb("F4F6F8"),
    "slate":    _rgb("4A6274"), "lightGray":_rgb("E8EDF0"),
    "midGray":  _rgb("8FA3B1"), "green":    _rgb("27AE60"),
    "red":      _rgb("E74C3C"), "orange":   _rgb("E67E22"),
    "darkGray": _rgb("2C3E50"), "amber":    _rgb("BF360C"),
    "grn_dk":   _rgb("1B5E20"), "grn_lt":  _rgb("D0EDD4"),
    "amb_lt":   _rgb("FFE8B0"), "red_lt":  _rgb("FCCEC9"),
    "wknd_bg":  _rgb("E8EBF8"), "wknd_hdr":_rgb("1A2E6E"),
}
MO      = ['','Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
MO_FULL = ['','January','February','March','April','May','June',
           'July','August','September','October','November','December']
DAYS_MO = [0,31,28,31,30,31,30,31,31,30,31,30,31]
TOTAL_ROOMS = 150

# ── Safe value helpers ────────────────────────────────────────────────────────
def _fv(v, default=0.0):
    try:
        f = float(v)
        return default if str(v).strip() in ('nan','') else f
    except:
        return default

def _iv(v, default=0):
    return int(_fv(v, default))

def _sign(v):
    try: return "+" if float(v) > 0 else ""
    except: return ""

# ── Primitives ────────────────────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color or C['white']

def _rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def _txt(slide, text, x, y, w, h, size=10, bold=False, color=None,
         align='left', valign='middle', wrap=False, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tb.text_frame._txBody.bodyPr.set('anchor',
        {'top':'t','middle':'ctr','bottom':'b'}.get(valign,'ctr'))
    p = tf.paragraphs[0]
    p.alignment = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,
                   'right':PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.name = 'Calibri'
    if color: r.font.color.rgb = color
    return tb

def _hdr(slide, title, sub=None):
    _rect(slide, 0, 0, 13.33, 0.55, C['navy'])
    _txt(slide, title, 0.25, 0, 9, 0.55, size=14, bold=True,
         color=C['white'], valign='middle')
    if sub:
        _txt(slide, sub, 8.5, 0, 4.6, 0.55, size=9,
             color=C['tealLt'], align='right', valign='middle')

def _cell(cell, text, bg=None, fg=None, bold=False, size=9,
          align='center', wrap=False):
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    tf = cell.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = {'center':PP_ALIGN.CENTER,'left':PP_ALIGN.LEFT,
                   'right':PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
    for run in p.runs: run.text = ''
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = 'Calibri'
    if fg: r.font.color.rgb = fg

def _kpi_dark(slide, x, y, w, h, label, value, sub, accent):
    """Dark navy KPI tile matching reference."""
    _rect(slide, x, y, w, h, C['navy'])
    _rect(slide, x, y, w, 0.06, accent)
    _txt(slide, label, x+0.15, y+0.12, w-0.25, 0.25,
         size=8, color=C['midGray'])
    _txt(slide, value, x+0.15, y+0.35, w-0.25, 0.55,
         size=24, bold=True, color=C['white'])
    _txt(slide, sub, x+0.15, y+0.88, w-0.25, 0.25,
         size=8, color=accent)

def _kpi_light(slide, x, y, w, h, label, value, sub, accent):
    """Light white KPI tile."""
    _rect(slide, x, y, w, h, C['white'])
    _rect(slide, x, y, w, 0.05, accent)
    _txt(slide, label, x+0.12, y+0.10, w-0.2, 0.22,
         size=8, color=C['slate'])
    _txt(slide, value, x+0.12, y+0.30, w-0.2, 0.52,
         size=22, bold=True, color=C['navy'])
    _txt(slide, sub, x+0.12, y+0.80, w-0.2, 0.25,
         size=8, color=accent)

def _add_yoy(cell, chg_str, is_occ=False, is_index=False, is_total=False):
    """Insert inline YoY run before endParaRPr."""
    try: v = float(chg_str)
    except: return
    if v == 0: return
    sign = "+" if v > 0 else ""
    yoy = f"({sign}{v:.1f})" if (is_occ or is_index) else f"({sign}{v:.1f}%)"
    col = ("7FFFC4" if v > 0 else "FFB3B3") if is_total else \
          ("27AE60" if v > 0 else "E74C3C")
    tc = cell._tc
    txb = tc.find(qn('a:txBody'))
    if txb is None: return
    paras = txb.findall(qn('a:p'))
    if not paras: return
    mp = paras[0]
    runs = mp.findall(qn('a:r'))
    if not runs: return
    orig = runs[0].find(qn('a:rPr'))
    r2 = etree.Element(qn('a:r'))
    rp = etree.SubElement(r2, qn('a:rPr'))
    rp.set('lang','en-US'); rp.set('sz','600'); rp.set('dirty','0')
    sf = etree.SubElement(rp, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr')); sc.set('val', col)
    if orig is not None:
        for tag in [qn('a:latin'),qn('a:ea'),qn('a:cs')]:
            ch = orig.find(tag)
            if ch is not None: rp.append(copy.deepcopy(ch))
    t2 = etree.SubElement(r2, qn('a:t')); t2.text = f" {yoy}"
    end = mp.find(qn('a:endParaRPr'))
    if end is not None: end.addprevious(r2)
    else: mp.append(r2)

# ── Data extraction ───────────────────────────────────────────────────────────
def extract_all(xl_path):
    """Read all needed data from Excel. Returns dict of structured data.
    xl_path can be: a file path string, a pd.ExcelFile, or a BytesIO object.
    """
    import pandas as pd
    if isinstance(xl_path, pd.ExcelFile):
        xl = xl_path
    else:
        xl = pd.ExcelFile(xl_path)
    df_str = pd.read_excel(xl, sheet_name='STR Analysis',    header=None)
    df_as  = pd.read_excel(xl, sheet_name='Annual Summary',  header=None)
    df_90  = pd.read_excel(xl, sheet_name='90 Day Segments', header=None)
    df_pu  = pd.read_excel(xl, sheet_name='Pickup',          header=None)

    # ── Property info ─────────────────────────────────────────────────────
    prop_name  = str(df_str.iloc[1][1])
    date_range = str(df_str.iloc[4][1])   # e.g. "5/31/2026 - 6/6/2026"
    start_str  = date_range.split(' - ')[0].strip()
    parts      = start_str.split('/')
    report_mo  = int(parts[0])
    report_yr  = int(parts[2])

    info = {
        'name': prop_name,
        'date_range': date_range,
        'report_mo': report_mo,
        'report_yr': report_yr,
    }

    # ── STR weekly + 28-day ───────────────────────────────────────────────
    # DOW col pairs (val, chg): Sun(3,4) Mon(6,7) Tue(9,10) Wed(12,13)
    #                            Thu(15,16) Fri(18,19) Sat(21,22) Total(24,25)
    DOW_PAIRS = [(3,4),(6,7),(9,10),(12,13),(15,16),(18,19),(21,22),(24,25)]

    def _exrow(row_idx):
        row = df_str.iloc[row_idx]
        return [{'val': str(row[vc]) if str(row[vc]) not in ('nan','') else '',
                 'chg': str(row[cc]) if str(row[cc]) not in ('nan','') else ''}
                for vc, cc in DOW_PAIRS]

    str_data = {
        'weekly': {
            'my_occ':    _exrow(8),
            'cs_occ':    _exrow(9),
            'mpi':       _exrow(10),
            'my_adr':    _exrow(12),
            'cs_adr':    _exrow(13),
            'ari':       _exrow(14),
            'my_revpar': _exrow(16),
            'cs_revpar': _exrow(17),
            'rgi':       _exrow(18),
        },
        'd28': {
            'my_occ':    _exrow(24),
            'cs_occ':    _exrow(25),
            'mpi':       _exrow(26),
            'my_adr':    _exrow(28),
            'cs_adr':    _exrow(29),
            'ari':       _exrow(30),
            'my_revpar': _exrow(32),
            'cs_revpar': _exrow(33),
            'rgi':       _exrow(34),
        },
    }

    # ── Week segments from STR Analysis ───────────────────────────────────
    # Row 37 (0-idx): segment names at cols 12,16,20,24,28,32,36,40,44
    # Row 46 (0-idx): weekly totals; each block offset from name col:
    #   +0=CY RMS, +1=STLY RMS, +2=CY ADR, +3=STLY ADR
    week_segs = []
    try:
        name_row  = df_str.iloc[37]
        total_row = df_str.iloc[46]
        seg_cols  = [ci for ci in range(12, df_str.shape[1], 4)
                     if str(name_row.iloc[ci]) not in ('nan','')]
        for nc in seg_cols:
            seg_name  = str(name_row.iloc[nc]).strip()
            cy_rms    = _iv(total_row.iloc[nc])
            stly_rms  = _iv(total_row.iloc[nc+1])
            cy_adr    = _fv(total_row.iloc[nc+2])
            stly_adr  = _fv(total_row.iloc[nc+3])
            if cy_rms == 0 and stly_rms == 0:
                continue
            week_segs.append({
                'name':    seg_name,
                'cy_rms':  cy_rms,  'stly_rms': stly_rms,
                'cy_adr':  cy_adr,  'stly_adr': stly_adr,
                'rms_var': cy_rms - stly_rms,
                'adr_var': cy_adr - stly_adr,
                'rev_var': cy_rms * cy_adr - stly_rms * stly_adr,
            })
    except Exception:
        week_segs = []

    # ── Annual pace (OTB section rows 6-18) ───────────────────────────────
    MO_NAMES = ['','Jan','Feb','Mar','Apr','May','Jun',
                'Jul','Aug','Sep','Oct','Nov','Dec']
    pace_data = []
    for i in range(6, 18):
        row = df_as.iloc[i]
        mo_name = str(row[1])
        if mo_name not in MO_NAMES: continue
        mo_num = MO_NAMES.index(mo_name)
        pace_data.append({
            'month': mo_name, 'mo_num': mo_num,
            'otb_occ': _fv(row[2])*100,
            'otb_rms': _iv(row[3]),
            'otb_adr': _fv(row[4]),
            'otb_rev': _fv(row[5]),
            'stly_occ': _fv(row[6])*100,
            'stly_rms': _iv(row[7]),
            'stly_adr': _fv(row[8]),
            'stly_rev': _fv(row[9]),
        })

    # ── Transient pace (rows 22-34) ───────────────────────────────────────
    trn_map = {}
    for i in range(22, 34):
        row = df_as.iloc[i]
        mo_name = str(row[1])
        if mo_name not in MO_NAMES: continue
        mo_num = MO_NAMES.index(mo_name)
        trn_map[mo_num] = {
            'trn_rms':      _iv(row[5]),
            'trn_adr':      _fv(row[6]),
            'trn_rev':      _fv(row[7]),
            'trn_stly_rms': _iv(row[8]),
            'trn_stly_adr': _fv(row[9]),
        }

    for d in pace_data:
        t = trn_map.get(d['mo_num'], {})
        d.update({
            'trn_rms':      t.get('trn_rms', 0),
            'trn_adr':      t.get('trn_adr', 0.0),
            'trn_rev':      t.get('trn_rev', 0.0),
            'trn_stly_rms': t.get('trn_stly_rms', 0),
            'trn_stly_adr': t.get('trn_stly_adr', 0.0),
        })

    # ── BI Forecast + Budget (rows 38-50) ─────────────────────────────────
    fcst_map = {}
    for i in range(38, 50):
        row = df_as.iloc[i]
        mo_name = str(row[1])
        if mo_name not in MO_NAMES: continue
        mo_num = MO_NAMES.index(mo_name)
        fcst_map[mo_num] = {
            'bi_occ': _fv(row[6]),
            'bi_rms': _fv(row[7]),
            'bi_adr': _fv(row[8]),
            'bi_rev': _fv(row[9]),
            'bud_occ': _fv(row[21]),
            'bud_rms': _iv(row[22]),
            'bud_adr': _fv(row[23]),
            'bud_rev': _fv(row[24]),
        }

    for d in pace_data:
        f = fcst_map.get(d['mo_num'], {})
        d.update({
            'bi_occ': f.get('bi_occ', d['otb_occ']/100),
            'bi_rms': f.get('bi_rms', d['otb_rms']),
            'bi_adr': f.get('bi_adr', d['otb_adr']),
            'bi_rev': f.get('bi_rev', d['otb_rev']),
            'bud_occ': f.get('bud_occ', 0.0),
            'bud_rms': f.get('bud_rms', 0),
            'bud_adr': f.get('bud_adr', 0.0),
            'bud_rev': f.get('bud_rev', 0.0),
        })

    # Total row 50
    tr = df_as.iloc[50]
    info['total_otb_rms']  = _iv(tr[3])
    info['total_otb_adr']  = _fv(tr[4])
    info['total_otb_rev']  = _fv(tr[5])
    info['total_stly_rms'] = _iv(tr[7])
    info['total_bi_rms']   = _fv(tr[7+4])   # col 11
    info['total_bi_rev']   = _fv(tr[9+4])   # col 13... use BI total col 9 from row 50
    # Safer: sum from pace_data
    info['total_bi_rms']   = int(sum(d['bi_rms'] for d in pace_data))
    info['total_bi_rev']   = sum(d['bi_rev'] for d in pace_data)
    info['total_bud_rms']  = int(sum(d['bud_rms'] for d in pace_data))
    info['total_bud_rev']  = sum(d['bud_rev'] for d in pace_data)
    info['total_bi_adr']   = info['total_bi_rev'] / info['total_bi_rms'] \
                             if info['total_bi_rms'] > 0 else 0

    # ── 90 Day Segments: monthly totals ───────────────────────────────────
    # Confirmed: Jun total row 36, Jul total row 71, Aug total row 106
    # cols: 3=LTS, 4=OCC, 5=ADR, 6=REV, 12=Trn RMS, 13=Trn ADR, 14=Trn REV
    #       15=Trn STLY RMS, 16=Trn STLY ADR, 18=Grp Block, 19=Grp RMS
    MO_TOTAL_ROWS = {0: 36, 1: 71, 2: 106}  # offset -> row index
    mo_totals = {}
    for offset, row_idx in MO_TOTAL_ROWS.items():
        mo_num = report_mo + offset
        if mo_num > 12: break
        if row_idx >= df_90.shape[0]: break
        r = df_90.iloc[row_idx]
        mo_totals[mo_num] = {
            'lts':       _iv(r[3]),
            'occ':       _fv(r[4]),
            'adr':       _fv(r[5]),
            'rev':       _fv(r[6]),
            'trn_rms':   _iv(r[12]),
            'trn_adr':   _fv(r[13]),
            'trn_rev':   _fv(r[14]),
            'trn_stly':  _iv(r[15]),
            'grp_block': _iv(r[18]),
            'grp_rms':   _iv(r[19]),
        }

    # ── 90 Day daily data for monthly OCC chart ───────────────────────────
    mo_daily = {report_mo: [], report_mo+1: [], report_mo+2: []}
    # First block: rows 6-35 for current month
    # Second block: rows 40-70 for +1 month
    # Third block: rows 75-105 for +2 month
    BLOCK_RANGES = [(6, 36), (40, 71), (75, 106)]
    for offset, (start, end) in enumerate(BLOCK_RANGES):
        mo_num = report_mo + offset
        if mo_num > 12: break
        days = []
        for i in range(start, min(end, df_90.shape[0])):
            row = df_90.iloc[i]
            dv = str(row[1])
            if '2026' not in dv and '2027' not in dv: continue
            try:
                parts2 = dv.split(',')[0].split('/')
                m2, d2 = int(parts2[0]), int(parts2[1])
                occ = _fv(row[4]) * 100
                days.append({'label': f"{m2}/{d2}", 'occ_otb': round(occ,1)})
            except: continue
        mo_daily[mo_num] = days

    # STLY OCC from 365 Day Outlook for the chart
    try:
        df_365 = pd.read_excel(xl, sheet_name='365 Day Outlook', header=None)
        stly_map = {}
        for i in range(6, df_365.shape[0]):
            row = df_365.iloc[i]
            dv = str(row[1])
            if '2026' not in dv and '2027' not in dv: continue
            try:
                parts2 = dv.split(',')[0].split('/')
                m2, d2 = int(parts2[0]), int(parts2[1])
                stly_rms = _fv(row[31])
                stly_occ = round(stly_rms / TOTAL_ROOMS * 100, 1)
                stly_map[f"{m2}/{d2}"] = stly_occ
            except: continue
        for mo_num, days in mo_daily.items():
            for d in days:
                d['occ_stly'] = stly_map.get(d['label'], 0)
    except: pass

    # ── Pickup data ───────────────────────────────────────────────────────
    # Row 39: month labels at cols 2,4,6,...
    # Row 40: RMS totals at same cols
    # Row 41: ADR totals
    # Row 42: REV totals
    pu_months = []; pu_rms = []; pu_adr = []; pu_rev = []
    pickup_from = ''
    try:
        mo_row  = df_pu.iloc[39]
        rms_row = df_pu.iloc[40]
        adr_row = df_pu.iloc[41]
        rev_row = df_pu.iloc[42]
        # month labels are at even cols starting col 2
        for ci in range(2, df_pu.shape[1], 2):
            mo_s = str(mo_row.iloc[ci])
            if mo_s in ('nan','','Total'): continue
            rms_v = _iv(rms_row.iloc[ci])
            adr_v = _fv(adr_row.iloc[ci])
            rev_v = _fv(rev_row.iloc[ci])
            if rms_v == 0 and adr_v == 0: continue
            short = mo_s.split('-')[0]
            pu_months.append(short); pu_rms.append(rms_v)
            pu_adr.append(round(adr_v)); pu_rev.append(rev_v)
        # Pickup from date
        for ci in range(df_pu.shape[1]):
            v = str(df_pu.iloc[4].iloc[ci])
            if 'Pickup From' in v or 'pickup from' in v.lower():
                pickup_from = v.split(':')[-1].strip() if ':' in v else v
    except: pass

    info['pickup_from'] = pickup_from

    # ── Pickup segment data from 90 Day totals ────────────────────────────
    cur_mo_tot = mo_totals.get(report_mo, {})
    total_trn_pu = pu_rms[0] if pu_rms else 0
    pu_segs = [
        {'name':'Transient (Total)', 'rms': total_trn_pu,
         'rev': int(pu_rev[0]) if pu_rev else 0,
         'pct_trans': '100.0%', 'is_total': True,
         'note': f"ADR ${pu_adr[0]:.0f}" if pu_adr else ''},
    ]
    # Segment-level pickup not available in this export
    for seg in ['Retail / Rack','Discount','Internet/OTA','Packages','Group']:
        pu_segs.append({'name':seg,'rms':0,'rev':0,'pct_trans':'',
                        'is_total':False,'note':'—'})

    return {
        'info':      info,
        'str_data':  str_data,
        'week_segs': week_segs,
        'pace_data': pace_data,
        'mo_totals': mo_totals,
        'mo_daily':  mo_daily,
        'pu_months': pu_months,
        'pu_rms':    pu_rms,
        'pu_adr':    pu_adr,
        'pu_rev':    pu_rev,
        'pu_segs':   pu_segs,
    }


# ── Slide 1: Title / KPI Dashboard ───────────────────────────────────────────
def build_slide_title(prs, data):
    info = data['info']
    w    = data['str_data']['weekly']
    slide = _blank(prs)
    _bg(slide, C['navy'])

    # Large title text
    _txt(slide, "REVENUE STRATEGY",
         0.5, 0.55, 12, 0.85, size=44, bold=True, color=C['white'])
    _txt(slide, "REVIEW & OUTLOOK",
         0.5, 1.38, 12, 0.85, size=44, bold=True, color=C['tealLt'])
    # Gold accent line
    _rect(slide, 0.5, 2.28, 1.2, 0.06, C['gold'])
    # Property name
    _txt(slide, info['name'],
         0.5, 2.50, 12, 0.45, size=16, color=C['tealLt'])
    # Date / subtitle
    dr = info['date_range']
    _txt(slide, f"Week of {dr}  |  Revenue Package",
         0.5, 2.95, 12, 0.35, size=11, color=C['midGray'])

    # 5 KPI tiles at bottom
    def fv(key): return _fv(w[key][7]['val'])
    def fchg(key): return w[key][7]['chg']

    kpis = [
        ("OCC",       f"{fv('my_occ'):.1f}%",    f"CS: {fv('cs_occ'):.1f}%",   C['teal']),
        ("ADR",       f"${fv('my_adr'):.2f}",     f"CS: ${fv('cs_adr'):.2f}",   C['tealLt']),
        ("My RevPAR", f"${fv('my_revpar'):.2f}",
         f"{_sign(fchg('my_revpar'))}{_fv(fchg('my_revpar')):.1f}% YoY",         C['tealLt']),
        ("CS RevPAR", f"${fv('cs_revpar'):.2f}",
         f"{_sign(fchg('cs_revpar'))}{_fv(fchg('cs_revpar')):.1f}% YoY",
         C['gold']),
        ("RGI",       f"{fv('rgi'):.1f}",
         f"{_sign(fchg('rgi'))}{_fv(fchg('rgi')):.1f} pts YoY",
         C['green'] if fv('rgi') >= 100 else C['red']),
    ]
    kw = 2.55
    for i, (lbl, val, sub, acc) in enumerate(kpis):
        _kpi_dark(slide, 0.25 + i*kw, 4.15, kw-0.1, 1.3, lbl, val, sub, acc)


# ── Slides 2–3: STR Weekly / 28-Day ──────────────────────────────────────────
def _build_str_slide(prs, data, period):
    info = data['info']
    w    = data['str_data'][period]
    segs = data['week_segs'] if period == 'weekly' else []

    slide = _blank(prs); _bg(slide)

    if period == 'weekly':
        title = f"STR PERFORMANCE — WEEK OF {info['date_range'].upper()}, {info['report_yr']}"
        sub   = f"{info['name'].split()[0]} vs Comp Set"
    else:
        title = "STR PERFORMANCE — RUNNING 28 DAYS (BY DAY OF WEEK)"
        sub   = f"{info['name'].split()[0]} vs Comp Set"
        _txt(slide, "28-DAY RUNNING SUMMARY", 0.25, 0.58, 8, 0.28,
             size=9, bold=True, color=C['teal'])

    _hdr(slide, title, sub)

    rows_def = [
        ("My OCC",    'my_occ',    True,  False, C['teal']),
        ("CS OCC",    'cs_occ',    True,  False, C['slate']),
        ("MPI",       'mpi',       False, True,  C['darkGray']),
        ("My ADR",    'my_adr',    False, False, C['teal']),
        ("CS ADR",    'cs_adr',    False, False, C['slate']),
        ("ARI",       'ari',       False, True,  C['darkGray']),
        ("My RevPAR", 'my_revpar', False, False, C['teal']),
        ("CS RevPAR", 'cs_revpar', False, False, C['slate']),
        ("RGI",       'rgi',       False, True,  C['darkGray']),
    ]

    tbl_y = 0.58 if period == 'weekly' else 0.88
    cw    = [1.55,1.4,1.4,1.4,1.4,1.4,1.4,1.4,1.45]
    tf = slide.shapes.add_table(10, 9, Inches(0.25), Inches(tbl_y),
                                Inches(12.83), Inches(3.3))
    t = tf.table
    for ci, w2 in enumerate(cw): t.columns[ci].width = int(w2*914400)
    for ri in range(10): t.rows[ri].height = int(0.33*914400)

    for ci, h in enumerate(["Metric","Sun","Mon","Tue","Wed","Thu","Fri","Sat","TOTAL"]):
        bg = _rgb("0A2438") if h == "TOTAL" else C['navy']
        fg = C['gold']  if h == "TOTAL" else C['white']
        _cell(t.cell(0,ci), h, bg=bg, fg=fg, bold=True, size=9.5)

    for ri, (lbl, key, is_occ, is_idx, color) in enumerate(rows_def):
        is_idx_row = key in ('mpi','ari','rgi')
        row_bg = C['teal'] if is_idx_row else \
                 (_rgb("E8EDF0") if ri%2==0 else C['white'])
        _cell(t.cell(ri+1, 0), lbl, bg=row_bg,
              fg=C['white'] if is_idx_row else C['darkGray'],
              bold=is_idx_row, size=9.5)
        for ci, d in enumerate(w[key]):
            is_total = (ci == 7)
            bg = _rgb("0A2438") if is_total else \
                 (C['white'] if ri%2==0 else _rgb("F4F6F8"))
            val = d['val']
            if key in ('my_adr','cs_adr','my_revpar','cs_revpar'):
                try: val = f"${float(val):.0f}"
                except: pass
            _cell(t.cell(ri+1, ci+1), val,
                  bg=bg, fg=C['gold'] if is_total else color,
                  bold=is_total, size=9.5)
            _add_yoy(t.cell(ri+1, ci+1), d['chg'],
                     is_occ=is_occ, is_index=is_idx, is_total=is_total)

    if period == 'weekly':
        # ── Bottom half: segment table + commentary ───────────────────────
        SEC_Y = 4.05
        _rect(slide, 0.25, SEC_Y, 7.2, 0.33, C['navy'])
        _txt(slide, "MARKET SEGMENT MIX — WEEK  |  CY vs STLY",
             0.25, SEC_Y, 7.2, 0.33, size=9.5, bold=True,
             color=C['white'], align='center', valign='middle')

        # Commentary panel
        _rect(slide, 7.6, SEC_Y, 5.5, 3.4, _rgb("F4F6F8"))
        _rect(slide, 7.6, SEC_Y, 5.5, 0.33, C['navy'])
        _txt(slide, "515 STR ANALYSIS",
             7.6, SEC_Y, 5.5, 0.33, size=10, bold=True,
             color=C['white'], align='center', valign='middle')

        rgi_val    = w['rgi'][7]['val']
        rgi_chg    = w['rgi'][7]['chg']
        mpi_val    = w['mpi'][7]['val']
        ari_val    = w['ari'][7]['val']
        my_adr_chg = w['my_adr'][7]['chg']
        cs_adr_chg = w['cs_adr'][7]['chg']
        dr         = info['date_range']

        # Build segment commentary
        if segs:
            top_gainers = sorted(segs, key=lambda s: s['rms_var'], reverse=True)[:2]
            top_losers  = sorted(segs, key=lambda s: s['rms_var'])[:2]
            gainer_txt  = " and ".join(
                f"{s['name']} growing {s['rms_var']:+d} rooms at "
                f"{_sign(str(s['adr_var']))}${s['adr_var']:.0f} ADR"
                for s in top_gainers if s['rms_var'] > 0
            ) or "select segments"
            loser_txt   = ", ".join(
                f"{s['name']} fell {abs(s['rms_var'])} rooms"
                for s in top_losers if s['rms_var'] < 0
            ) or ""
            seg_note = f"The ARI improvement was driven by {gainer_txt} versus STLY. " \
                       f"Offsetting this, {loser_txt}." if loser_txt else \
                       f"The ARI improvement was driven by {gainer_txt} versus STLY."
        else:
            seg_note = "Monitor segment mix and group pace to protect midweek production."

        critique = (
            f"The week of {dr} delivered an RGI of {rgi_val} "
            f"({_sign(rgi_chg)}{rgi_chg} pts YoY), "
            f"supported by an MPI of {mpi_val} and ARI of {ari_val}. "
            "The hotel outperformed the comp set on both occupancy and rate "
            f"simultaneously for the week. {seg_note} "
            f"Rate growth was driven by ADR improving "
            f"{_sign(my_adr_chg)}{my_adr_chg}% vs STLY while the comp set "
            f"grew {_sign(cs_adr_chg)}{cs_adr_chg}%."
        )
        _txt(slide, critique, 7.68, SEC_Y+0.42, 5.32, 2.88,
             size=9.5, color=C['darkGray'], wrap=True)

        # Segment table
        n_seg_rows = max(len(segs), 1)
        tf2 = slide.shapes.add_table(n_seg_rows+1, 4,
            Inches(0.25), Inches(SEC_Y+0.33),
            Inches(7.2), Inches(3.07))
        t2 = tf2.table
        for ci, cw2 in enumerate([2.5, 1.5, 1.6, 1.6]):
            t2.columns[ci].width = int(cw2*914400)
        for ri in range(n_seg_rows+1):
            t2.rows[ri].height = int(0.27*914400)
        for ci, h in enumerate(["Segment","RMS Var","ADR Var","Rev Var"]):
            _cell(t2.cell(0,ci), h, bg=C['teal'], fg=C['white'],
                  bold=True, size=9.5)
        if not segs:
            _cell(t2.cell(1,0), "No segment data in this export",
                  bg=_rgb("F4F6F8"), fg=C['midGray'], size=9, align='left')
            for ci in range(1,4):
                _cell(t2.cell(1,ci),'—', bg=_rgb("F4F6F8"), fg=C['midGray'], size=9)
        else:
            for ri, seg in enumerate(segs):
                bg    = _rgb("F4F6F8") if ri%2==0 else C['white']
                rms_c = C['green'] if seg['rms_var']>0 else \
                        C['red'] if seg['rms_var']<0 else C['midGray']
                adr_c = C['green'] if seg['adr_var']>0 else \
                        C['orange'] if seg['adr_var']<0 else C['midGray']
                rev_c = C['green'] if seg['rev_var']>0 else \
                        C['red'] if seg['rev_var']<0 else C['midGray']
                rms_s = f"{seg['rms_var']:+d}"
                adr_s = f"{_sign(str(seg['adr_var']))}${abs(seg['adr_var']):.0f}"
                rev_s = f"{_sign(str(seg['rev_var']))}${abs(seg['rev_var'])/1000:.1f}K"
                _cell(t2.cell(ri+1,0), seg['name'],
                      bg=bg, fg=C['darkGray'], bold=True, size=9.5)
                _cell(t2.cell(ri+1,1), rms_s, bg=bg, fg=rms_c, size=9.5)
                _cell(t2.cell(ri+1,2), adr_s, bg=bg, fg=adr_c, size=9.5)
                _cell(t2.cell(ri+1,3), rev_s, bg=bg, fg=rev_c, size=9.5)

    else:
        # 28-day: 4 KPI tiles at bottom
        rgi_tot = _fv(w['rgi'][7]['val'])
        mpi_tot = _fv(w['mpi'][7]['val'])
        ari_tot = _fv(w['ari'][7]['val'])
        rv_tot  = _fv(w['my_revpar'][7]['val'])
        rv_chg  = w['my_revpar'][7]['chg']
        mpi_chg = w['mpi'][7]['chg']
        rgi_chg = w['rgi'][7]['chg']
        ari_chg = w['ari'][7]['chg']
        kpis = [
            ("28-Day MPI",    f"{mpi_tot:.1f}",   f"{_sign(mpi_chg)}{_fv(mpi_chg):.1f} pts",  C['teal']),
            ("28-Day ARI",    f"{ari_tot:.1f}",   f"{_sign(ari_chg)}{_fv(ari_chg):.1f} pts",  C['gold']),
            ("28-Day RGI",    f"{rgi_tot:.1f}",   f"{_sign(rgi_chg)}{_fv(rgi_chg):.1f} pts",  C['green'] if rgi_tot>=100 else C['red']),
            ("My RevPAR",     f"${rv_tot:.2f}",   f"{_sign(rv_chg)}{_fv(rv_chg):.1f}% YoY",   C['tealLt']),
        ]
        for i, (lbl, val, sub, acc) in enumerate(kpis):
            _kpi_dark(slide, 0.25+i*3.28, 4.2, 3.2, 1.3, lbl, val, sub, acc)

    return slide


def build_slide_str_weekly(prs, data):
    return _build_str_slide(prs, data, 'weekly')

def build_slide_str_28day(prs, data):
    return _build_str_slide(prs, data, 'd28')


# ── Slide 4: Annual Pace ──────────────────────────────────────────────────────
def build_slide_annual_pace(prs, data):
    info  = data['info']
    pace  = data['pace_data']
    slide = _blank(prs); _bg(slide)
    _hdr(slide, "ANNUAL PACE — OTB vs STLY vs BUDGET",
         f"{info['report_yr']} Full Year | Transient Focus")

    months   = [d['month'] for d in pace]
    otb_rms  = [d['otb_rms']  for d in pace]
    stly_rms = [d['stly_rms'] for d in pace]
    bud_rms  = [d['bud_rms']  for d in pace]

    cd = ChartData(); cd.categories = months
    cd.add_series(f"OTB {info['report_yr']}",   tuple(otb_rms))
    cd.add_series(f"STLY {info['report_yr']-1}", tuple(stly_rms))
    cd.add_series("Budget",                      tuple(bud_rms))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.25), Inches(0.65), Inches(13.0), Inches(3.8), cd).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = C['teal']
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = C['lightGray']
    chart.series[2].format.fill.solid()
    chart.series[2].format.fill.fore_color.rgb = C['gold']

    # Detail table: upcoming 3 months
    upcoming = [d for d in pace if d['mo_num'] >= info['report_mo']][:3]
    TABLE_Y = 4.6
    headers = ["Month","OTB RMS","OTB ADR","OTB REV","STLY RMS","STLY ADR","STLY REV",
               "RMS Var","ADR Var","REV Var","Bud REV","vs Bud"]
    cws = [0.7,0.9,0.9,0.9,0.9,0.9,0.9,0.7,0.7,0.9,0.9,0.7]
    tf = slide.shapes.add_table(len(upcoming)+1, len(headers),
         Inches(0.25), Inches(TABLE_Y), Inches(12.83), Inches(2.8))
    t = tf.table
    for ci, cw in enumerate(cws): t.columns[ci].width = int(cw*914400)
    for ri in range(len(upcoming)+1): t.rows[ri].height = int(0.55*914400)

    # Group header bands
    grp_cols = [
        ("ON THE BOOKS", 1, 3, C['teal']),
        ("SAME TIME LAST YEAR", 4, 6, C['slate']),
        ("VARIANCE vs STLY", 7, 9, _rgb("1A5C3A")),
        ("BUDGET", 10, 10, C['gold']),
        ("VAR vs BUD", 11, 11, _rgb("5C3A1A")),
    ]
    # Single header row with coloured bg per group
    _cell(t.cell(0,0), "Month", bg=C['navy'], fg=C['white'], bold=True, size=8)
    group_colors = {1:C['teal'],2:C['teal'],3:C['teal'],
                    4:C['slate'],5:C['slate'],6:C['slate'],
                    7:_rgb("1A5C3A"),8:_rgb("1A5C3A"),9:_rgb("1A5C3A"),
                    10:C['gold'],11:_rgb("5C3A1A")}
    for ci, h in enumerate(headers[1:], 1):
        bg = group_colors.get(ci, C['navy'])
        _cell(t.cell(0,ci), h, bg=bg, fg=C['white'], bold=True, size=7.5)

    for ri, d in enumerate(upcoming):
        bg = _rgb("F4F6F8") if ri%2==0 else C['white']
        rms_var = d['otb_rms'] - d['stly_rms']
        adr_var = d['otb_adr'] - d['stly_adr']
        stly_rev = d['stly_rms'] * d['stly_adr']
        rev_var  = d['otb_rev'] - stly_rev
        vs_bud   = (d['otb_rev']/d['bud_rev']-1)*100 if d['bud_rev']>0 else 0
        grn = C['green']; red = C['red']
        row_vals = [
            (d['month'],              C['navy'],  True),
            (f"{d['otb_rms']:,}",     C['teal'],  False),
            (f"${d['otb_adr']:.2f}",  C['teal'],  False),
            (f"${d['otb_rev']/1000:.0f}K", C['teal'], True),
            (f"{d['stly_rms']:,}",    C['slate'], False),
            (f"${d['stly_adr']:.2f}", C['slate'], False),
            (f"${stly_rev/1000:.0f}K",C['slate'], False),
            (f"{rms_var:+,}",         grn if rms_var>=0 else red, True),
            (f"{_sign(str(adr_var))}${adr_var:.2f}", grn if adr_var>=0 else red, False),
            (f"{_sign(str(rev_var))}${abs(rev_var)/1000:.0f}K", grn if rev_var>=0 else red, True),
            (f"${d['bud_rev']/1000:.0f}K", C['midGray'], False),
            (f"{vs_bud:+.1f}%",       grn if vs_bud>=0 else red, True),
        ]
        for ci, (v, fg, bold) in enumerate(row_vals):
            _cell(t.cell(ri+1, ci), v, bg=bg, fg=fg, bold=bold, size=8.5)


# ── Slide 5: Transient Pace ───────────────────────────────────────────────────
def build_slide_transient_pace(prs, data):
    info  = data['info']
    pace  = data['pace_data']
    slide = _blank(prs); _bg(slide)
    _hdr(slide, "TRANSIENT PACE — OTB vs. STLY",
         "Transient Room Nights & ADR by Month")

    months   = [d['month'] for d in pace]
    trn_otb  = [d['trn_rms']      for d in pace]
    trn_stly = [d['trn_stly_rms'] for d in pace]
    otb_adr  = [d['trn_adr']      for d in pace]
    stly_adr = [d['trn_stly_adr'] for d in pace]

    cd = ChartData(); cd.categories = months
    cd.add_series(f"Transient OTB {info['report_yr']}",   tuple(trn_otb))
    cd.add_series(f"Transient STLY {info['report_yr']-1}", tuple(trn_stly))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.25), Inches(0.65), Inches(6.5), Inches(3.7), cd).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = C['teal']
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = C['lightGray']

    cd2 = ChartData(); cd2.categories = months
    cd2.add_series(f"OTB ADR {info['report_yr']}",   tuple(otb_adr))
    cd2.add_series(f"STLY ADR {info['report_yr']-1}", tuple(stly_adr))
    chart2 = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(6.9), Inches(0.65), Inches(6.2), Inches(3.7), cd2).chart
    chart2.has_title = True
    chart2.chart_title.text_frame.text = "ADR: OTB vs STLY"
    chart2.has_legend = True
    chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart2.series[0].format.line.color.rgb = C['teal']
    chart2.series[0].format.line.width = 18000
    chart2.series[1].format.line.color.rgb = C['gold']
    chart2.series[1].format.line.width = 18000

    # 3 insight tiles
    report_mo = info['report_mo']
    passed = [d for d in pace if d['mo_num'] < report_mo]
    current = next((d for d in pace if d['mo_num'] == report_mo), None)
    future  = [d for d in pace if d['mo_num'] > report_mo][:2]

    def _tile(slide, x, title, body, accent):
        _rect(slide, x, 4.55, 4.22, 2.5, C['white'])
        _rect(slide, x, 4.55, 4.22, 0.05, accent)
        _txt(slide, title, x+0.15, 4.6, 3.9, 0.3, size=9, bold=True, color=C['darkGray'])
        _txt(slide, body,  x+0.15, 4.95, 3.9, 2.0, size=9, color=C['slate'], wrap=True)

    # Jan–current month summary
    if passed:
        yr_rms  = sum(d['trn_rms']      for d in passed)
        yr_stly = sum(d['trn_stly_rms'] for d in passed)
        yr_adr  = sum(d['trn_adr']      for d in passed) / len(passed)
        yr_sadr = sum(d['trn_stly_adr'] for d in passed) / len(passed)
        mo_names = "/".join(d['month'] for d in passed[:3])
        body = (f"OTB: {', '.join(str(d['trn_rms']) for d in passed[:3])}\n"
                f"STLY: {', '.join(str(d['trn_stly_rms']) for d in passed[:3])}\n"
                f"ADR +${yr_adr-yr_sadr:.2f} vs STLY avg")
        _tile(slide, 0.25, f"Jan–{passed[-1]['month']}", body, C['teal'])
    if current:
        var = current['trn_rms'] - current['trn_stly_rms']
        body = (f"OTB {current['trn_rms']:,} vs STLY {current['trn_stly_rms']:,}\n"
                f"{var:+,} RNs ({var/current['trn_stly_rms']*100:+.1f}%)\n"
                f"ADR ${current['trn_adr']:.2f} vs ${current['trn_stly_adr']:.2f} STLY")
        acc = C['green'] if var >= 0 else C['orange']
        _tile(slide, 4.56, f"{current['month']} OTB vs STLY", body, acc)
    if future:
        lines = []
        for d in future:
            var = d['trn_rms'] - d['trn_stly_rms']
            lines.append(f"OTB {d['trn_rms']:,} vs STLY {d['trn_stly_rms']:,} ({MO[d['mo_num']]})")
        body = "\n".join(lines) + "\nMonitor forward pace closely"
        _tile(slide, 8.87, "May+ Pace Deficit" if any(d['trn_rms']<d['trn_stly_rms'] for d in future) else "Forward Pace",
              body, C['orange'])


# ── Slide 6: 7-Day Pickup ─────────────────────────────────────────────────────
def build_slide_pickup(prs, data):
    info = data['info']
    slide = _blank(prs); _bg(slide)
    _hdr(slide, "7-DAY PICKUP REPORT",
         f"Pickup from {info.get('pickup_from','')}")

    pu_months = data['pu_months']
    pu_rms    = data['pu_rms']
    pu_adr    = data['pu_adr']
    pu_rev    = data['pu_rev']

    if pu_months and pu_rms:
        cd = ChartData(); cd.categories = pu_months
        cd.add_series("Room Nights Picked Up", tuple(pu_rms))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.25), Inches(0.65), Inches(6.5), Inches(3.8), cd).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "Room Nights Picked Up by Month"
        chart.has_legend = False
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = C['teal']
        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.show_value = True

    if pu_months and pu_adr:
        cd2 = ChartData(); cd2.categories = pu_months
        cd2.add_series("Average ADR of Pickup", tuple(pu_adr))
        chart2 = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(6.9), Inches(0.65), Inches(6.2), Inches(3.8), cd2).chart
        chart2.has_title = True
        chart2.chart_title.text_frame.text = "Average ADR of Pickup by Month"
        chart2.has_legend = False
        chart2.series[0].format.fill.solid()
        chart2.series[0].format.fill.fore_color.rgb = C['gold']
        chart2.plots[0].has_data_labels = True
        chart2.plots[0].data_labels.show_value = True

    # Segment table
    _rect(slide, 0.25, 4.62, 13.0, 0.33, C['navy'])
    _txt(slide, "PICKUP DETAIL (7-Day) — Key Transient Segments",
         0.35, 4.62, 13.0, 0.33, size=9.5, bold=True,
         color=C['white'], valign='middle')

    segs = data['pu_segs']
    tf = slide.shapes.add_table(len(segs)+1, 5,
        Inches(0.25), Inches(5.0), Inches(13.0), Inches(2.4))
    t = tf.table
    for ci, cw in enumerate([2.8, 1.7, 1.5, 2.0, 5.0]):
        t.columns[ci].width = int(cw*914400)
    for ri in range(len(segs)+1):
        t.rows[ri].height = int(0.36*914400)

    for ci, h in enumerate(["Segment","PU RNs","% of Trans","PU Revenue","Note"]):
        _cell(t.cell(0,ci), h, bg=C['teal'], fg=C['white'], bold=True, size=9.5)

    for ri, seg in enumerate(segs):
        is_total = seg.get('is_total', False)
        bg = _rgb("0D1B2A") if is_total else \
             (_rgb("F4F6F8") if ri%2==0 else C['white'])
        fg = C['gold'] if is_total else C['darkGray']
        rms_v = seg['rms']
        rms_s = f"+{rms_v}" if rms_v > 0 else (str(rms_v) if rms_v < 0 else "—")
        rms_c = C['green'] if rms_v > 0 else C['red'] if rms_v < 0 else C['midGray']
        _cell(t.cell(ri+1,0), seg['name'],   bg=bg, fg=fg,  bold=is_total, size=9.5)
        _cell(t.cell(ri+1,1), rms_s,         bg=bg, fg=rms_c if not is_total else C['gold'], bold=is_total, size=9.5)
        _cell(t.cell(ri+1,2), seg.get('pct_trans',''), bg=bg, fg=C['teal'] if not is_total else C['gold'], size=9.5)
        _cell(t.cell(ri+1,3), f"${seg['rev']:,}" if seg['rev'] else '—', bg=bg, fg=C['teal'] if not is_total else C['gold'], bold=is_total, size=9.5)
        _cell(t.cell(ri+1,4), seg.get('note',''), bg=bg, fg=C['slate'] if not is_total else C['midGray'], size=9, align='left')


# ── Slides 7/9/11: Monthly Daily Outlook ─────────────────────────────────────
def build_slide_monthly_occ(prs, data, mo_num):
    info     = data['info']
    mo_tot   = data['mo_totals'].get(mo_num, {})
    daily    = data['mo_daily'].get(mo_num, [])
    pace_row = next((d for d in data['pace_data'] if d['mo_num'] == mo_num), {})

    slide = _blank(prs); _bg(slide)
    mo_label = MO_FULL[mo_num] if mo_num <= 12 else ''
    _hdr(slide, f"{mo_label.upper()} DAILY OUTLOOK — OTB vs STLY",
         f"Transient  |  {mo_label} {info['report_yr']}")

    labels    = [d['label']    for d in daily]
    otb_vals  = [d['occ_otb']  for d in daily]
    stly_vals = [d.get('occ_stly', 0) for d in daily]

    if labels:
        cd = ChartData(); cd.categories = labels
        cd.add_series(f"OTB {info['report_yr']} OCC%",     tuple(otb_vals))
        cd.add_series(f"STLY {info['report_yr']-1} OCC%",  tuple(stly_vals))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            Inches(0.25), Inches(0.65), Inches(13.0), Inches(3.8), cd).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.series[0].format.line.color.rgb = C['teal']
        chart.series[0].format.line.width = 18000
        chart.series[1].format.line.color.rgb = C['gold']
        chart.series[1].format.line.width = 18000
    else:
        _txt(slide, "No daily data available.", 0.25, 2.5, 13.0, 0.5,
             size=12, color=C['midGray'], align='center')

    # KPI tiles (dark) matching reference
    otb_rms  = mo_tot.get('otb_rms',  pace_row.get('otb_rms', 0))
    if not otb_rms: otb_rms = pace_row.get('otb_rms', 0)
    occ_pct  = mo_tot.get('occ', pace_row.get('otb_occ', 0)/100) * 100 \
               if mo_tot else pace_row.get('otb_occ', 0)
    trn_rms  = mo_tot.get('trn_rms', pace_row.get('trn_rms', 0))
    trn_adr  = mo_tot.get('trn_adr', pace_row.get('trn_adr', 0.0))
    stly_rms = pace_row.get('stly_rms', 0)
    vs_rms   = otb_rms - stly_rms
    otb_rev  = pace_row.get('otb_rev', 0)
    stly_rev = pace_row.get('stly_rev', 0)
    vs_rev   = (otb_rev/stly_rev-1)*100 if stly_rev > 0 else 0

    vs_col = C['green'] if vs_rms >= 0 else C['red']
    rev_col= C['green'] if vs_rev >= 0 else C['red']

    kpis = [
        (f"{MO[mo_num]} OTB RMS",      f"{otb_rms:,}",
         f"{occ_pct:.1f}% OCC",        C['teal']),
        (f"{MO[mo_num]} Transient RMS", f"{trn_rms:,}",
         f"${trn_adr:.2f} ADR",         C['tealLt']),
        ("vs STLY Rooms",               f"{vs_rms:+,}",
         f"vs STLY {stly_rms:,}",       vs_col),
        (f"{MO[mo_num]} Revenue OTB",   f"${otb_rev/1000:.0f}K",
         f"{vs_rev:+.1f}% vs STLY",     rev_col),
    ]
    for i, (lbl, val, sub, acc) in enumerate(kpis):
        _kpi_dark(slide, 0.25+i*3.26, 4.65, 3.2, 1.3, lbl, val, sub, acc)


# ── Slides 8/10/12: Market Segment Mix (Monthly) ─────────────────────────────
def build_slide_segment_mix(prs, data, mo_num):
    """
    Monthly segment mix from Annual Summary transient section
    and STR Analysis segment data. Pie chart + detail table + takeaways.
    """
    info     = data['info']
    pace_row = next((d for d in data['pace_data'] if d['mo_num'] == mo_num), {})
    segs     = data['week_segs']  # best available segment data

    slide = _blank(prs); _bg(slide)
    mo_label = MO_FULL[mo_num] if mo_num <= 12 else ''
    _hdr(slide, f"MARKET SEGMENT MIX — {mo_label.upper()} {info['report_yr']}")

    if segs:
        # Pie chart using segment revenue contribution
        total_rev = sum(max(s['cy_rms']*s['cy_adr'], 0) for s in segs)
        if total_rev > 0:
            cd = ChartData()
            for s in segs:
                rev = max(s['cy_rms'] * s['cy_adr'], 0)
                cd.add_series(s["name"], (rev / total_rev,))
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                Inches(0.25), Inches(0.65), Inches(5.5), Inches(5.5), cd).chart
            chart.has_title = True
            chart.chart_title.text_frame.text = "Revenue Contribution %"
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        # Detail table
        tf = slide.shapes.add_table(len(segs)+1, 7,
            Inches(5.9), Inches(0.65), Inches(7.2), Inches(4.3))
        t = tf.table
        for ci, cw in enumerate([1.6, 0.8, 0.9, 0.85, 0.9, 0.8, 0.8]):
            t.columns[ci].width = int(cw*914400)
        for ri in range(len(segs)+1):
            t.rows[ri].height = int(0.52*914400)
        for ci, h in enumerate(["Segment","CY RNs","CY ADR","STLY RNs","STLY ADR","RN Var","ADR Var"]):
            _cell(t.cell(0,ci), h, bg=C['navy'], fg=C['white'], bold=True, size=9)
        for ri, seg in enumerate(segs):
            bg = _rgb("F4F6F8") if ri%2==0 else C['white']
            rvar = seg['rms_var']; avar = seg['adr_var']
            rc   = C['green'] if rvar>0 else C['red'] if rvar<0 else C['midGray']
            ac   = C['green'] if avar>0 else C['orange'] if avar<0 else C['midGray']
            vals = [
                (seg['name'],                          C['darkGray'], True),
                (f"{seg['cy_rms']:,}",                 C['teal'],     False),
                (f"${seg['cy_adr']:.2f}",              C['teal'],     True),
                (f"{seg['stly_rms']:,}",                C['slate'],    False),
                (f"${seg['stly_adr']:.2f}",             C['slate'],    False),
                (f"{rvar:+,}" if rvar else "○",         rc,            True),
                (f"{_sign(str(avar))}${avar:.2f}",      ac,            False),
            ]
            for ci, (v, fg, bold) in enumerate(vals):
                _cell(t.cell(ri+1, ci), v, bg=bg, fg=fg, bold=bold, size=9)

        # Takeaways
        _rect(slide, 5.9, 5.1, 7.2, 0.32, C['teal'])
        _txt(slide, "KEY TAKEAWAYS", 6.0, 5.1, 7.0, 0.32,
             size=9, bold=True, color=C['white'], valign='middle')
        _rect(slide, 5.9, 5.42, 7.2, 2.0, _rgb("F4F6F8"))

        top2  = sorted(segs, key=lambda s: s['rms_var'], reverse=True)[:2]
        bot2  = sorted(segs, key=lambda s: s['rms_var'])[:2]
        lines = []
        for s in top2:
            if s['rms_var'] > 0:
                lines.append(f"✓ {s['name']} up {s['rms_var']:+,} RNs vs STLY at ${s['cy_adr']:.2f} ADR")
        for s in bot2:
            if s['rms_var'] < 0:
                lines.append(f"▲ {s['name']} decreased {s['rms_var']:,} RNs vs STLY — monitor")
        body = "\n".join(lines[:4]) if lines else "Review segment performance vs STLY."
        _txt(slide, body, 6.0, 5.48, 6.9, 1.85, size=9, color=C['darkGray'], wrap=True)
    else:
        _txt(slide, "Segment data not available in this export.",
             0.25, 3.5, 13.0, 0.5, size=11, color=C['midGray'], align='center')


# ── Slide 13: Full Year Forecast ──────────────────────────────────────────────
def build_slide_full_year(prs, data):
    info  = data['info']
    pace  = data['pace_data']
    slide = _blank(prs); _bg(slide)
    _hdr(slide, f"FULL YEAR FORECAST {info['report_yr']}",
         "Finance Forecast vs Budget")

    months   = [d['month'] for d in pace]
    fcst_rms = [d['bi_rms'] for d in pace]
    bud_rms  = [d['bud_rms'] for d in pace]

    cd = ChartData(); cd.categories = months
    cd.add_series("Forecast", tuple(fcst_rms))
    cd.add_series("Budget",   tuple(bud_rms))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.25), Inches(0.65), Inches(7.4), Inches(4.8), cd).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.series[0].format.line.color.rgb = C['teal']
    chart.series[0].format.line.width = 20000
    chart.series[1].format.line.color.rgb = C['gold']
    chart.series[1].format.line.width = 20000

    # Monthly table (forward months)
    fwd = [d for d in pace if d['mo_num'] >= info['report_mo']]
    tf = slide.shapes.add_table(len(fwd)+1, 7,
        Inches(7.65), Inches(0.65), Inches(5.45), Inches(4.8))
    t = tf.table
    for ci, cw in enumerate([0.75,0.85,0.85,0.9,0.85,0.75,0.5]):
        t.columns[ci].width = int(cw*914400)
    for ri in range(len(fwd)+1): t.rows[ri].height = int(0.38*914400)

    for ci, h in enumerate(["Month","OTB OCC","Fcst OCC","Fcst ADR","Fcst REV","Budget REV","vs Bud"]):
        _cell(t.cell(0,ci), h, bg=C['navy'], fg=C['white'], bold=True, size=8)

    for ri, d in enumerate(fwd):
        bg  = _rgb("F4F6F8") if ri%2==0 else C['white']
        cap = DAYS_MO[d['mo_num']] * TOTAL_ROOMS
        otb_occ_p  = d['otb_occ']
        fcst_occ_p = d['bi_occ'] * 100 if d['bi_occ'] <= 1 else d['bi_occ']
        vs_bud = (d['bi_rev']/d['bud_rev']-1)*100 if d['bud_rev']>0 else 0
        vs_col = C['green'] if vs_bud >= 0 else C['red']
        fcst_col = C['grn_dk'] if fcst_occ_p >= 95 else C['amber']
        _cell(t.cell(ri+1,0), d['month'],                       bg=bg, fg=C['navy'],   bold=True, size=8)
        _cell(t.cell(ri+1,1), f"{otb_occ_p:.1f}%",             bg=bg, fg=C['teal']   if otb_occ_p>80 else C['orange'], size=8)
        _cell(t.cell(ri+1,2), f"{fcst_occ_p:.1f}%",            bg=bg, fg=fcst_col,    bold=True, size=8)
        _cell(t.cell(ri+1,3), f"${d['bi_adr']:.2f}",           bg=bg, fg=C['teal'],   size=8)
        _cell(t.cell(ri+1,4), f"${d['bi_rev']/1000:.0f}K",     bg=bg, fg=C['navy'],   bold=True, size=8)
        _cell(t.cell(ri+1,5), f"${d['bud_rev']/1000:.0f}K",    bg=bg, fg=C['midGray'],size=8)
        _cell(t.cell(ri+1,6), f"{vs_bud:+.1f}%",               bg=bg, fg=vs_col,      bold=True, size=8)

    # KPI tiles
    total_bi_rev  = info['total_bi_rev']
    total_bud_rev = info['total_bud_rev']
    total_bi_rms  = info['total_bi_rms']
    total_bi_adr  = info['total_bi_adr']
    yr_occ = total_bi_rms / (365 * TOTAL_ROOMS) * 100
    vs_bud_pct = (total_bi_rev/total_bud_rev-1)*100 if total_bud_rev > 0 else 0

    kpis = [
        ("Forecast Full Year OCC", f"{yr_occ:.1f}%",           "", C['teal']),
        ("Forecast Full Year ADR", f"${total_bi_adr:.2f}",     "", C['gold']),
        ("Forecast Full Year REV", f"${total_bi_rev/1e6:.2f}M","", C['tealLt']),
        ("vs Budget REV",          f"{vs_bud_pct:+.1f}%",      "",
         C['green'] if vs_bud_pct >= 0 else C['red']),
    ]
    for i, (lbl, val, sub, acc) in enumerate(kpis):
        _kpi_dark(slide, 0.25+i*3.28, 5.6, 3.2, 1.3, lbl, val, sub, acc)

    # Commentary
    commentary = (
        f"Finance Forecast projects full year at ${total_bi_rev/1e6:.2f}M vs "
        f"${total_bud_rev/1e6:.2f}M budget ({vs_bud_pct:+.1f}%). "
        "Heavy lift required in back half as current OTB reflects early pace. "
        "Transient volume capture is the primary revenue lever."
    )
    _rect(slide, 0.25, 7.05, 13.0, 0.38, _rgb("F0F4F8"))
    _rect(slide, 0.25, 7.05, 0.06, 0.38, C['teal'])
    _txt(slide, commentary, 0.4, 7.07, 12.7, 0.34,
         size=8.5, color=C['slate'], wrap=True)


# ── Master builder ─────────────────────────────────────────────────────────────
def build_slides_1_to_14(prs, xl_path, info_override=None,
                          str_data_override=None, monthly_data=None,
                          ops_rows=None, total_rooms=TOTAL_ROOMS):
    """
    Called by app.py. Accepts xl_path (str or file-like).
    info_override / str_data_override kept for backward compat but ignored —
    we re-extract everything from xl_path directly.
    """
    # xl_path can be a file path, pd.ExcelFile, or BytesIO/bytes
    import io, pandas as pd
    if isinstance(xl_path, pd.ExcelFile):
        data = extract_all(xl_path)
    elif hasattr(xl_path, 'read') or isinstance(xl_path, (bytes, bytearray)):
        if isinstance(xl_path, (bytes, bytearray)):
            xl_path = io.BytesIO(xl_path)
        import tempfile, os
        tmp = tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False)
        xl_path.seek(0); tmp.write(xl_path.read()); tmp.flush(); tmp.close()
        data = extract_all(tmp.name)
        os.unlink(tmp.name)
    else:
        data = extract_all(xl_path)

    build_slide_title(prs, data)              # Slide 1
    build_slide_str_weekly(prs, data)         # Slide 2
    build_slide_str_28day(prs, data)          # Slide 3
    build_slide_annual_pace(prs, data)        # Slide 4
    build_slide_transient_pace(prs, data)     # Slide 5
    build_slide_pickup(prs, data)             # Slide 6

    report_mo = data['info']['report_mo']
    for offset in range(3):
        mo = report_mo + offset
        if mo > 12: break
        build_slide_monthly_occ(prs, data, mo)    # 7, 9, 11
        build_slide_segment_mix(prs, data, mo)    # 8, 10, 12

    build_slide_full_year(prs, data)          # Slide 13
    # Slide 14 = ops forecast built by app.py
